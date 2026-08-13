"""노션 1교시의 「실제 노션 화면」을 원본 자료에서 그대로 꺼낸다.

    python 노션화면추출.py

무엇이며 왜 있는가
    강사가 실제 노션 화면을 캡처해 보냈고, 그 화면을 도형으로 1:1 재현해 만든 것이
    최초의 노션 pptx 다. 덱을 코드 빌드로 옮기면서 그 재현을 손그림 목업으로
    바꿔치기한 적이 있는데, 그러면 화면과 설명이 어긋난다. **다시 그리지 않는다.**
    원본을 고해상도로 뽑아 화면 영역만 잘라 쓴다.

    원본은 저장소에 없다(pptx 는 산출물이라 .gitignore 로 막혀 있다).
    커밋 626e20e 에 마지막 손제작본이 남아 있어 거기서 꺼낸다.

    결과는 특강/그림/노션/*.png — 이건 자산이라 커밋한다.
    (한 번 뽑아 두면 다시 돌릴 일이 없다. 원본이 바뀌지 않기 때문이다.)

필요한 것
    PowerPoint (COM 으로 고해상도 내보내기) · python-pptx · Pillow
"""

from __future__ import annotations

import subprocess
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parents[1]
그림 = REPO / "강의" / "그림" / "노션"

원본커밋 = "626e20e"
원본경로 = "특강/2일차/이론/노션_자산화_강의자료.pptx"

# 잘라 쓸 화면 — (원본 장, 이름, 손지정 사각형 또는 None)
#   None 이면 창 표시줄의 신호등 원 3개를 찾아 창틀을 자동으로 잡는다.
#   사각형은 (x, y, w, h) 를 슬라이드 크기 대비 비율로 적는다 — 창틀이 없는
#   메뉴·팝오버용이다.
화면들: list[tuple[int, str, tuple[float, float, float, float] | None]] = [
    (4,  "완성_홈",          None),
    (5,  "완성_카드",        None),
    (7,  "만들기_01_페이지", None),
    (8,  "만들기_02_블록",   None),
    (10, "만들기_04_표메뉴", (0.585, 0.216, 0.355, 0.385)),
    (11, "만들기_05_유형",   (0.600, 0.228, 0.325, 0.312)),
    (11, "만들기_05_옵션",   (0.600, 0.553, 0.325, 0.186)),
    (12, "만들기_06_카드",   None),
    (13, "만들기_07_템플릿", None),
    (14, "만들기_08_본문",   None),
    (15, "만들기_09_코드",   None),
    (16, "만들기_10_체크",   None),
    (17, "만들기_11_갤러리", None),
    (18, "만들기_12_보드",   (0.556, 0.234, 0.408, 0.345)),
    (18, "만들기_12_남은것", (0.556, 0.601, 0.408, 0.299)),
    (19, "만들기_13_공유",   (0.615, 0.246, 0.310, 0.320)),
    (21, "완성_최종",        None),
]

내보내기 = """param([string]$Src, [string]$Out, [int]$W, [int]$H)
$ErrorActionPreference = "Stop"
if (-not (Test-Path $Out)) { New-Item -ItemType Directory -Path $Out | Out-Null }
$app = New-Object -ComObject PowerPoint.Application
$pres = $app.Presentations.Open($Src, -1, 0, -1)
for ($i = 1; $i -le $pres.Slides.Count; $i++) {
    $pres.Slides.Item($i).Export((Join-Path $Out ("{0:d2}.png" -f $i)), "PNG", $W, $H)
}
$pres.Close(); $app.Quit()
Write-Output ("exported " + $pres.Slides.Count)
"""


def 신호등창틀(slide, W, H):
    """창 표시줄의 작은 원 3개를 찾고, 그것을 품은 가장 큰 도형을 창틀로 본다."""
    무리: dict[float, list] = {}
    for sh in slide.shapes:
        if sh.left is None or sh.width is None:
            continue
        x, y, w, h = sh.left / W, sh.top / H, sh.width / W, sh.height / H
        글자 = sh.has_text_frame and sh.text_frame.text.strip()
        if 0.002 < w < 0.012 and 0.004 < h < 0.020 and not 글자:
            무리.setdefault(round(y, 3), []).append((x, y))
    등 = next((v for v in 무리.values() if len(v) >= 3), None)
    if not 등:
        return None
    lx, ly = min(t[0] for t in 등), min(t[1] for t in 등)
    후보 = []
    for sh in slide.shapes:
        if sh.left is None or sh.width is None:
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            continue
        x, y, w, h = sh.left / W, sh.top / H, sh.width / W, sh.height / H
        if x <= lx + 1e-3 and y <= ly + 1e-3 and x + w >= lx and y + h >= ly + 5e-3 and h > 0.05:
            후보.append((h, x, y, w, h))
    if not 후보:
        return None
    _, x, y, w, h = max(후보)
    return (x, y, w, h)


def main() -> int:
    try:
        from PIL import Image
        from pptx import Presentation
    except ModuleNotFoundError as e:
        print(f"{e.name} 가 없습니다 — 제작/ 환경에서 돌리세요", file=sys.stderr)
        return 1

    임시 = Path(tempfile.mkdtemp(prefix="노션원본_"))
    pptx = 임시 / "원본.pptx"
    r = subprocess.run(["git", "cat-file", "-p", f"{원본커밋}:{원본경로}"],
                       cwd=REPO, capture_output=True)
    if r.returncode != 0 or not r.stdout:
        print(f"원본을 못 꺼냈습니다 — {원본커밋}:{원본경로}", file=sys.stderr)
        print(r.stderr.decode("utf-8", "replace")[:300], file=sys.stderr)
        return 1
    pptx.write_bytes(r.stdout)

    ps1 = 임시 / "내보내기.ps1"
    ps1.write_text(내보내기, encoding="ascii")
    고해상 = 임시 / "고해상"
    r = subprocess.run(["powershell", "-ExecutionPolicy", "Bypass", "-File", str(ps1),
                        "-Src", str(pptx), "-Out", str(고해상), "-W", "3200", "-H", "1800"],
                       capture_output=True, text=True)
    if not any(고해상.glob("*.png")):
        print("PowerPoint 내보내기가 실패했습니다", file=sys.stderr)
        print((r.stdout + r.stderr)[-600:], file=sys.stderr)
        return 1

    prs = Presentation(str(pptx))
    W, H = prs.slide_width, prs.slide_height
    장 = list(prs.slides)

    그림.mkdir(parents=True, exist_ok=True)
    여백 = 0.006
    for 번호, 이름, 사각 in 화면들:
        틀 = 사각 or 신호등창틀(장[번호 - 1], W, H)
        if not 틀:
            print(f"  [빠짐] {번호}장 창틀을 못 찾았습니다 — {이름}", file=sys.stderr)
            return 1
        x, y, w, h = 틀
        im = Image.open(고해상 / f"{번호:02d}.png").convert("RGB")
        iw, ih = im.size
        box = (max(0, int((x - 여백) * iw)), max(0, int((y - 여백) * ih)),
               min(iw, int((x + w + 여백) * iw)), min(ih, int((y + h + 여백) * ih)))
        나온것 = im.crop(box)
        나온것.save(그림 / f"{이름}.png")
        print(f"  {이름:<18} {나온것.width}x{나온것.height}  (원본 {번호}장)")

    print(f"\n  {len(화면들)}장 →  {그림}")
    return 0


if __name__ == "__main__":
    sys.exit(main())

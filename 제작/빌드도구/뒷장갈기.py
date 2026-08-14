# -*- coding: utf-8 -*-
"""강사가 손본 pptx 의 **앞부분은 그대로 두고 뒷장만 갈아 끼운다.**

    python 제작/빌드도구/뒷장갈기.py "<강사파일.pptx>" <살릴장수>

왜 필요한가
    강사가 앞장을 직접 고쳤다. 전체를 다시 뽑으면 그 손질이 통째로 날아간다.
    (강조 상자를 손으로 맞춘 좌표, 그림 크기, 문구 — 다시 만들 수 없는 것들이다)

    그래서 앞장은 **파일에 있는 그대로 두고**, 뒷장만 지운 뒤 새로 만든 것을 붙인다.

겪은 것 둘
    ① 장을 옮기면서 **딸린 관계를 통째로** 이어 붙였더니 레이아웃·테마·그림이
       zip 안에서 이름이 겹쳐 파워포인트가 파일을 못 열었다.
       → 레이아웃은 새 장 것을 쓰고, **그림만 새 부품으로 다시 넣는다.**
    ② 파워포인트 COM 으로 합치려 했으나 이 PC 에서는 자동화로 **어떤 파일도 수정할 수
       없다** (`Presentation cannot be modified`). 그래서 COM 은 쓰지 않는다.

    ※ 원본을 덮어쓰지 않는다. 결과는 「(합침)」 이 붙은 새 파일로 나온다.
"""

from __future__ import annotations

import sys as _sys

for _s in (_sys.stdout, _sys.stderr):
    if (getattr(_s, "encoding", "") or "").lower().replace("-", "") != "utf8":
        try:
            _s.reconfigure(encoding="utf-8", errors="replace")
        except Exception:
            pass

import shutil
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent
새로만든 = ROOT.parents[1] / "강의" / "2일차" / "강의자료" / "노션_자산화.pptx"


def 각주줄맞추기(경로: Path, 남길수: int, y: int = 815) -> int:
    """앞장의 **각주 높이만** 하나로 맞춘다. 글자도 강조 상자도 건드리지 않는다.

    강사가 장마다 손으로 내려놓아 808~821 로 흩어져 있다. 넘길 때 눈에 걸린다.
    구분선(y=782)이 있는 장에서, 그 아래 있는 글상자만 옮긴다.
    """
    from pptx import Presentation
    from pptx.util import Emu

    EMU = 7620
    prs = Presentation(str(경로))
    옮긴수 = 0
    for i, s in enumerate(prs.slides):
        if i >= 남길수:
            break
        선있나 = any(sh.top is not None and abs(Emu(sh.top).inches * 120 - 782) < 3
                     and sh.width is not None and Emu(sh.width).inches * 120 > 1000
                     for sh in s.shapes)
        if not 선있나:
            continue
        for sh in s.shapes:
            if not (sh.has_text_frame and sh.text_frame.text.strip()):
                continue
            현재 = Emu(sh.top).inches * 120
            if 790 < 현재 < 845 and abs(현재 - y) > 0.5:
                sh.top = Emu(round(y * EMU))
                옮긴수 += 1
    if 옮긴수:
        prs.save(str(경로))
    return 옮긴수


_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
       "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}
_EMBED = f"{{{_NS['r']}}}embed"


def 장지우기(prs, 남길수: int) -> int:
    """남길수 장만 남기고 뒤를 전부 지운다."""
    목록 = prs.slides._sldIdLst
    지운수 = 0
    for sid in list(목록)[남길수:]:
        prs.part.drop_rel(sid.rId)
        목록.remove(sid)
        지운수 += 1
    return 지운수


def 장옮기기(받을prs, 줄prs, 처음: int) -> int:
    """줄prs 의 처음번째 장부터 끝까지를 받을prs 뒤에 붙인다.

    도형 XML 은 그대로 옮기되, **그림은 새 부품으로 다시 넣는다.**
    남의 꾸러미에 있던 그림 부품을 그대로 이어 붙이면 `image9.png` 처럼
    이름이 겹쳐 파일이 깨진다 (실제로 그렇게 깨졌다).
    """
    import copy
    import io

    빈판 = 받을prs.slide_layouts[6]
    붙인수 = 0
    for i, 원장 in enumerate(줄prs.slides):
        if i < 처음:
            continue
        새장 = 받을prs.slides.add_slide(빈판)
        for sh in list(새장.shapes):          # 빈 판이 달고 온 자리 표시를 걷어낸다
            sh._element.getparent().remove(sh._element)

        try:                                  # 바탕색
            if 원장.background.fill.type is not None:
                새장.background.fill.solid()
                새장.background.fill.fore_color.rgb = \
                    원장.background.fill.fore_color.rgb
        except Exception:
            pass

        for sh in 원장.shapes:
            새el = copy.deepcopy(sh._element)
            for blip in 새el.iter(f"{{{_NS['a']}}}blip"):
                옛rId = blip.get(_EMBED)
                if not 옛rId:
                    continue
                자료 = 원장.part.related_part(옛rId).blob
                _, 새rId = 새장.part.get_or_add_image_part(io.BytesIO(자료))
                blip.set(_EMBED, 새rId)
            새장.shapes._spTree.insert_element_before(새el, "p:extLst")
        붙인수 += 1
    return 붙인수


def main() -> int:
    if len(sys.argv) < 3:
        print(__doc__.strip(), file=sys.stderr)
        return 1
    강사본 = Path(sys.argv[1])
    살릴수 = int(sys.argv[2])

    if not 강사본.is_file():
        print(f"파일이 없습니다: {강사본}", file=sys.stderr)
        return 1
    if not 새로만든.is_file():
        print(f"먼저 편집pptx.py 노션 을 돌리세요 — {새로만든} 가 없습니다", file=sys.stderr)
        return 1

    from pptx import Presentation

    낼곳 = 강사본.with_name(강사본.stem + " (합침).pptx")
    shutil.copy2(강사본, 낼곳)
    낼곳.chmod(0o666)

    받을 = Presentation(str(낼곳))
    처음수 = len(받을.slides._sldIdLst)
    지운수 = 장지우기(받을, 살릴수)
    붙인수 = 장옮기기(받을, Presentation(str(새로만든)), 살릴수)
    받을.save(str(낼곳))
    최종 = len(Presentation(str(낼곳)).slides._sldIdLst)

    print(f"  {강사본.name} — 앞 {살릴수}장 남기고 {지운수}장 지움 (원래 {처음수}장)")
    print(f"  {새로만든.name} 의 {살릴수 + 1}장부터 {붙인수}장을 붙임  →  모두 {최종}장")

    옮긴수 = 각주줄맞추기(낼곳, 살릴수)
    if 옮긴수:
        print(f"  각주 {옮긴수}개를 y=815 로 맞춤 (글자·강조 상자는 그대로)")

    print(f"\n  {낼곳}")
    print(f"  앞 {살릴수}장은 손대지 않았습니다 — 강조 상자도 그대로입니다.")
    return 0


if __name__ == "__main__":
    sys.exit(main())

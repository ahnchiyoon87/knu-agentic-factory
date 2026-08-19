"""손수정본의 정렬·여백을 새로 뽑은 pptx 에 그대로 이식한다.

    uv run --with python-pptx 손좌표적용.py

왜 있나
    강사가 파워포인트에서 상자를 직접 옮겨 정렬·여백을 손질한 판(수정본)이 있다.
    빌더는 브라우저 실측 좌표로 다시 그리므로, 다시 뽑을 때마다 그 손질이
    되돌아간다 — 실제로 11장에서 그랬다.

    그래서 **글자가 수정본과 완전히 같은 장**에 한해, 도형의 자리·크기를
    수정본 값으로 덮어쓴다. 글자가 달라진 장(이번에 고친 장)은 건드리지 않는다 —
    옛 좌표를 새 글자에 씌우면 상자가 글자를 못 담는다.

기준 좌표는 `제작/덱빌드/손좌표.json` 이다. `--추출` 로 수정본에서 다시 뽑는다.
수정본 원본(D:\\다운로드\\수정본)이 없어도 json 만 있으면 적용은 된다.
"""
from __future__ import annotations

# ── 한글 윈도우(cp949)에서 출력이 깨져 죽는 것을 막는다 ──────────────────
import sys as _sys
for _s in (_sys.stdout, _sys.stderr):
    if (getattr(_s, "encoding", "") or "").lower().replace("-", "") != "utf8":
        try:
            _s.reconfigure(encoding="utf-8", errors="replace")
        except Exception:
            pass
# ─────────────────────────────────────────────────────────────────────────

import argparse
import hashlib
import json
import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]          # 제작/
좌표파일 = ROOT / "덱빌드" / "손좌표.json"
수정본폴더 = Path(r"D:\다운로드\수정본")

대상 = [
    ("2일차", ROOT.parent / "강의" / "2일차" / "강의자료" / "슬라이드.pptx",
     수정본폴더 / "2일차 슬라이드_수정본.pptx"),
    ("3일차", ROOT.parent / "강의" / "3일차" / "강의자료" / "슬라이드.pptx",
     수정본폴더 / "3일차 슬라이드_수정본.pptx"),
]


def 장열쇠(slide) -> tuple[str, list[str]]:
    """장을 알아보는 열쇠 — 글자 있는 도형의 텍스트 목록(위치 순)."""
    글 = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = "\n".join(p.text for p in sh.text_frame.paragraphs).strip()
            if t:
                글.append((t, int(sh.top or 0), int(sh.left or 0)))
    글.sort(key=lambda x: (x[1], x[2]))
    본문 = [t for t, _, _ in 글]
    지문 = hashlib.sha256("\x1f".join(본문).encode("utf-8")).hexdigest()[:16]
    return 지문, 본문


def 추출() -> None:
    """수정본에서 「글자 지문 → 도형별 좌표」를 뽑아 json 으로 남긴다."""
    표 = {}
    for 이름, _새, 수정본 in 대상:
        if not 수정본.is_file():
            sys.exit(f"수정본이 없습니다: {수정본}")
        prs = Presentation(str(수정본))
        for slide in prs.slides:
            지문, 본문 = 장열쇠(slide)
            도형들 = []
            for sh in slide.shapes:
                if not sh.has_text_frame:
                    continue
                t = "\n".join(p.text for p in sh.text_frame.paragraphs).strip()
                if t:
                    도형들.append({"글": t, "l": int(sh.left or 0), "t": int(sh.top or 0),
                                   "w": int(sh.width or 0), "h": int(sh.height or 0)})
            표[f"{이름}:{지문}"] = {"제목": (본문[0].splitlines()[0] if 본문 else ""),
                                    "도형": 도형들}
    좌표파일.write_text(json.dumps(표, ensure_ascii=False, indent=1), encoding="utf-8")
    print(f"손좌표.json — 장 {len(표)}개 좌표를 뽑았습니다")


def 적용() -> None:
    if not 좌표파일.is_file():
        sys.exit("손좌표.json 이 없습니다 — 먼저 --추출 을 돌리세요.")
    표 = json.loads(좌표파일.read_text(encoding="utf-8"))
    for 이름, 새경로, _ in 대상:
        prs = Presentation(str(새경로))
        맞춘장, 만진도형 = 0, 0
        for slide in prs.slides:
            지문, _본문 = 장열쇠(slide)
            기록 = 표.get(f"{이름}:{지문}")
            if not 기록:
                continue                        # 글자가 달라진 장 — 손대지 않는다
            자리 = {}
            for d in 기록["도형"]:
                자리.setdefault(d["글"], []).append(d)
            바뀜 = False
            for sh in slide.shapes:
                if not sh.has_text_frame:
                    continue
                t = "\n".join(p.text for p in sh.text_frame.paragraphs).strip()
                큐 = 자리.get(t)
                if not 큐:
                    continue
                d = 큐.pop(0)
                if (int(sh.left or 0), int(sh.top or 0),
                        int(sh.width or 0), int(sh.height or 0)) != (d["l"], d["t"], d["w"], d["h"]):
                    sh.left, sh.top, sh.width, sh.height = d["l"], d["t"], d["w"], d["h"]
                    만진도형 += 1
                    바뀜 = True
            if 바뀜:
                맞춘장 += 1
        prs.save(str(새경로))
        print(f"{이름} — 손좌표를 이식한 장 {맞춘장}개 · 도형 {만진도형}개")


if __name__ == "__main__":
    ap = argparse.ArgumentParser(description="손수정본 좌표 이식")
    ap.add_argument("--추출", action="store_true", help="수정본 pptx 에서 좌표를 다시 뽑는다")
    args = ap.parse_args()
    if args.추출:
        추출()
    적용()

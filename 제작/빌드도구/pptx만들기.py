"""슬라이드 PNG 를 pptx 한 벌로 묶는다 — **모양이 절대 안 틀어지는 백업본**.

    python pptx만들기.py        →  강의/<일차>/강의자료/슬라이드_이미지본.pptx

강의장에서 넘기는 것은 이게 아니다
    강의장 덱은 `편집pptx.py` 가 만드는 **`슬라이드.pptx`** 다. 글자가 진짜 텍스트 상자라
    강의 직전에 파워포인트에서 문구를 고칠 수 있다. 드라이브에도 그쪽이 실린다.

    여기서 나오는 것은 장마다 PNG 한 장을 붙인 것이라 **한 자도 못 고친다.**
    대신 어느 PC 에서 열어도 화면과 100% 같다. 발표 PC 가 낯설 때의 보험이다.

    **pptx 를 손으로 고치지 말 것.** 다음에 다시 뽑으면 사라진다.
    글자를 고치려면 `build_*.py` 를 고치고 render 한 뒤 다시 뽑는다.
"""

from __future__ import annotations

import sys as _sys
for _s in (_sys.stdout, _sys.stderr):
    if (getattr(_s, "encoding", "") or "").lower().replace("-", "") != "utf8":
        try:
            _s.reconfigure(encoding="utf-8", errors="replace")
        except Exception:
            pass

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent          # 제작/빌드도구/
REPO = ROOT.parents[1]                          # 경남대특강/ (저장소 루트)
강의 = REPO / "강의"                             # pptx 가 놓일 곳: 강의/{일차}/강의자료/
PNG = REPO / "제작" / "산출물" / "슬라이드"        # render.py 가 찍어 둔 낱장


def main() -> int:
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ModuleNotFoundError:
        print("python-pptx 가 없습니다 —  pip install python-pptx", file=sys.stderr)
        return 1
    try:
        from PIL import Image
    except ModuleNotFoundError:
        print("Pillow 가 없습니다 —  pip install Pillow", file=sys.stderr)
        return 1

    # (PNG 폴더, 만들 pptx)
    # 노션은 여기서 만들지 않는다 — `편집pptx.py 노션` 이 **요소별로 고칠 수 있는**
    # pptx 를 만든다. 여기서 다시 뽑으면 통 이미지로 덮어써서 편집이 사라진다.
    #
    # 2·3일차도 같은 이유로 **이름을 나눈다.** 전에는 여기서도 `슬라이드.pptx` 에 썼는데,
    # `편집pptx.py` 가 쓰는 파일과 이름이 같아 **나중에 돌린 쪽이 이겼다.**
    # 이미지본을 나중에 돌리면 글자를 한 자도 못 고치는 덱이 드라이브에 실리고,
    # 드라이브 빌더는 둘을 구분하지 못한다. 강의장에서 넘기는 것은 `슬라이드.pptx`(편집본)이고,
    # 여기서 나오는 것은 **모양이 절대 안 틀어지는 백업본**이다.
    묶음 = [("2일차", "2일차", "슬라이드_이미지본.pptx"),
            ("3일차", "3일차", "슬라이드_이미지본.pptx")]

    for 일차, 폴더, 파일 in 묶음:
        src = PNG / 폴더
        pngs = sorted(src.glob("*.png"), key=lambda p: int(p.stem))
        if not pngs:
            print(f"  [빠짐] {src} 에 PNG 가 없습니다", file=sys.stderr)
            return 1

        w, h = Image.open(pngs[0]).size
        prs = Presentation()
        # 슬라이드 크기를 PNG 비율에 정확히 맞춘다 — 안 맞추면 여백이 생기거나 잘린다
        prs.slide_width = Emu(int(12192000))                      # 16:9 기준 폭
        prs.slide_height = Emu(int(12192000 * h / w))

        빈레이아웃 = prs.slide_layouts[6]                          # 완전 빈 장
        for p in pngs:
            s = prs.slides.add_slide(빈레이아웃)
            s.shapes.add_picture(str(p), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)

        dst = 강의 / 일차 / "강의자료" / 파일
        dst.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(dst))
        print(f"  만듦  {일차}/{dst.name}  ({len(pngs)}장 · {dst.stat().st_size // 1024}KB)")

    print("\n  ※ pptx 를 손으로 고치지 마세요. 다시 뽑으면 사라집니다.")
    print("     글자를 고치려면 —  제작/덱빌드/build_*.py  →  render.py  →  이 스크립트")
    return 0


if __name__ == "__main__":
    sys.exit(main())

"""노션 화면에서 **가리킬 대상의 좌표를 계산해** 준다 — 눈대중을 없애기 위한 것.

    python 노션좌표뽑기.py                 # 정해 둔 대상 전부
    python 노션좌표뽑기.py 만들기_05_유형   # 한 화면만

무엇이며 왜 있는가
    강조 상자를 손으로 적으면 매번 크거나 작거나 서로 겹친다. 실제로 그랬다.
    원본 pptx 에는 도형마다 정확한 위치가 들어 있으므로, **가리킬 것을 글자로 지목**하면
    그 도형들을 감싸는 최소 사각형을 계산할 수 있다. 그 값을 build_노션.py 가 쓴다.

    결과는 `제작/덱빌드/노션좌표.py` 에 파이썬 딕셔너리로 떨어진다.
    좌표를 고치고 싶으면 **여기 대상 목록을 고치고 다시 돌린다.** 숫자를 손으로 만지지 않는다.
"""

from __future__ import annotations

import subprocess
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parents[1]
결과파일 = REPO / "제작" / "덱빌드" / "노션좌표.py"

원본커밋 = "626e20e"
원본경로 = "특강/2일차/이론/노션_자산화_강의자료.pptx"

# 화면을 잘라낸 사각형 — 노션화면추출.py 와 같은 값이어야 한다
여백 = 0.006
손지정 = {
    "만들기_04_표메뉴": (10, (0.585, 0.216, 0.355, 0.385)),
    "만들기_05_유형":   (11, (0.600, 0.228, 0.325, 0.312)),
    "만들기_05_옵션":   (11, (0.600, 0.553, 0.325, 0.186)),
    "만들기_12_보드":   (18, (0.556, 0.234, 0.408, 0.345)),
    "만들기_12_남은것": (18, (0.556, 0.601, 0.408, 0.299)),
    "만들기_13_공유":   (19, (0.615, 0.246, 0.310, 0.320)),
}
자동 = {"완성_홈": 4, "완성_카드": 5, "만들기_01_페이지": 7, "만들기_02_블록": 8,
        "만들기_06_카드": 12, "만들기_07_템플릿": 13, "만들기_08_본문": 14,
        "만들기_09_코드": 15, "만들기_10_체크": 16, "만들기_11_갤러리": 17,
        "완성_최종": 21}

# ── 가리킬 대상 — 화면마다 (이름, [도형에 적힌 글자, ...]) ────────────────────
#    글자로 지목하면 그 도형들을 감싸는 최소 사각형이 상자가 된다.
#    같은 줄의 아이콘까지 감싸려면 아이콘 글자도 같이 적는다.
대상 = {
    "완성_홈": [
        ("사이드바",   ["📚 내 학습 노트", "📄 8/11 이상 감지", "📄 8/12 …"]),
        ("뷰탭",       ["🖼 갤러리", "▦ 표", "▥ 분야별", "🔁 남은 것 2"]),
        ("카드들",     ["📊 8/11 이상 감지", "설명가능", "🔧 8/13 파이프라인", "설명못함"]),
    ],
    "완성_카드": [
        ("속성",       ["◐ 분야", "AI 실습", "◐ 이해도", "설명가능"]),
        ("개념하나",   ["z-score가 재는 것", "평균에서 몇 칸 떨어졌는지를 재는 값. 단위가 달라도 같은 잣대로 비교할 수 있다."]),
        ("확인할것",   ["다음에 확인할 것", "윈도 크기 정하는 기준", "z가 3인 이유"]),
    ],
    "만들기_01_페이지": [
        ("새페이지",   ["＋ 새 페이지"]),
        ("제목자리",   ["내 학습 노트"]),
        ("아이콘",     ["📚"]),
    ],
    "만들기_02_블록": [
        ("콜아웃",     ["📌", "내가 이해한 걸 먼저 쓴다."]),
        ("제목2",      ["수업 기록"]),
    ],
    "만들기_04_표메뉴": [
        ("슬래시표",   ["/표"]),
        ("표보기",     ["▤", "표 보기", "인라인 데이터베이스를 추가합니다"]),
        ("그냥표",     ["≡", "표", "간단한 표 블록"]),
    ],
    "만들기_05_유형": [
        ("선택",       ["◐", "선택 ✓"]),
        ("다중선택",   ["≔", "다중 선택"]),
    ],
    "만들기_05_옵션": [
        ("옵션셋",     ["설명못함", "대충앎", "설명가능"]),
    ],
    "만들기_06_카드": [
        ("새로만들기", ["＋ 새로 만들기"]),
        ("이름칸",     ["📊 8/11 이상 감지"]),
        ("속성칸",     ["AI 실습", "설명못함"]),
    ],
    "만들기_07_템플릿": [
        ("본문",       ["개념 이름", "여기에 내 말로"]),
        ("확인할것",   ["다음에 확인할 것", "여기에"]),
    ],
    "만들기_08_본문": [
        ("개념1",      ["개념 이름", "여기에 내 말로 한 줄"]),
    ],
    "만들기_09_코드": [
        ("코드블록",   ["roll = s.rolling(window=60)"]),
        ("언어",       ["Python ▾"]),
        ("그림자리",   ["강의자료 캡처"]),
    ],
    "만들기_10_체크": [
        ("할일",       ["다음에 확인할 것", "못 들은 것 하나", "다시 볼 것 하나", "해결한 것 하나"]),
    ],
    "만들기_11_갤러리": [
        ("더하기",     ["＋"]),
        ("갤러리탭",   ["🖼 갤러리"]),
        ("카드하나",   ["📊 8/11 이상 감지", "설명못함"]),
    ],
    "만들기_12_보드": [
        ("분야별탭",   ["▥ 분야별"]),
        ("칸머리",     ["AI 실습"]),
    ],
    "만들기_12_남은것": [
        ("남은것탭",   ["🔁 남은 것"]),
        ("필터",       ["필터: 이해도 ≠ 설명가능"]),
        ("개수줄",     ["2개 · 여기가 비면 끝난 것"]),
    ],
    "만들기_13_공유": [
        ("게시탭",     ["게시"]),
        ("웹에게시",   ["웹에 게시"]),
        ("링크",       ["notion.site/내-학습-노트-a4f9c2…"]),
    ],
}


def 신호등창틀(slide, W, H):
    무리: dict[float, list] = {}
    for sh in slide.shapes:
        if sh.left is None or sh.width is None:
            continue
        x, y, w, h = sh.left / W, sh.top / H, sh.width / W, sh.height / H
        글자 = sh.has_text_frame and sh.text_frame.text.strip()
        if 0.002 < w < 0.012 and 0.004 < h < 0.020 and not 글자:
            무리.setdefault(round(y, 3), []).append((x, y))
    등 = next((v for v in 무리.values() if len(v) >= 3), None)
    if not 등:
        return None
    lx, ly = min(t[0] for t in 등), min(t[1] for t in 등)
    후보 = []
    for sh in slide.shapes:
        if sh.left is None or sh.width is None:
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            continue
        x, y, w, h = sh.left / W, sh.top / H, sh.width / W, sh.height / H
        if x <= lx + 1e-3 and y <= ly + 1e-3 and x + w >= lx and y + h >= ly + 5e-3 and h > .05:
            후보.append((h, x, y, w, h))
    return max(후보)[1:] if 후보 else None


def main() -> int:
    try:
        from pptx import Presentation
    except ModuleNotFoundError:
        print("python-pptx 가 없습니다 — 제작/ 환경에서 돌리세요", file=sys.stderr)
        return 1

    임시 = Path(tempfile.mkdtemp(prefix="노션좌표_"))
    pptx = 임시 / "원본.pptx"
    r = subprocess.run(["git", "cat-file", "-p", f"{원본커밋}:{원본경로}"],
                       cwd=REPO, capture_output=True)
    if r.returncode != 0 or not r.stdout:
        print("원본 pptx 를 못 꺼냈습니다", file=sys.stderr)
        return 1
    pptx.write_bytes(r.stdout)

    prs = Presentation(str(pptx))
    W, H = prs.slide_width, prs.slide_height
    장 = list(prs.slides)

    고른것 = sys.argv[1:] or list(대상)
    결과: dict[str, dict[str, tuple]] = {}
    빠짐 = []

    for 이름 in 고른것:
        if 이름 in 손지정:
            n, 틀 = 손지정[이름]
        else:
            n = 자동[이름]
            틀 = 신호등창틀(장[n - 1], W, H)
        cx, cy = 틀[0] - 여백, 틀[1] - 여백
        cw, ch = 틀[2] + 여백 * 2, 틀[3] + 여백 * 2

        칸 = {}
        for 별명, 글자들 in 대상[이름]:
            상자 = []
            for sh in 장[n - 1].shapes:
                if sh.left is None or sh.width is None or not sh.has_text_frame:
                    continue
                t = sh.text_frame.text.strip()
                # 여러 줄이 한 도형인 경우가 있다(코드 블록). 긴 글자만 앞부분으로도 맞춘다 —
                # 짧은 글자로 앞부분을 맞추면 「표」가 「표 보기」까지 잡아 상자가 커진다.
                맞음 = t in 글자들 or any(len(g) >= 10 and t.startswith(g) for g in 글자들)
                if not t or not 맞음:
                    continue
                x, y = sh.left / W, sh.top / H
                w, h = sh.width / W, sh.height / H
                # **잘라낸 화면 안**에 있는 도형만 쓴다. 같은 글자가 슬라이드 다른 곳에도
                # 있어서(설명문·다른 창) 함께 묶이면 상자가 화면 밖까지 커진다.
                if x < cx - 1e-4 or y < cy - 1e-4 or x + w > cx + cw + 1e-4 or y + h > cy + ch + 1e-4:
                    continue
                상자.append((x, y, x + w, y + h))
            if not 상자:
                빠짐.append(f"{이름}/{별명} — 글자를 못 찾음: {글자들}")
                continue
            x1 = min(b[0] for b in 상자)
            y1 = min(b[1] for b in 상자)
            x2 = max(b[2] for b in 상자)
            y2 = max(b[3] for b in 상자)
            상 = (round((x1 - cx) / cw * 100, 1), round((y1 - cy) / ch * 100, 1),
                  round((x2 - x1) / cw * 100, 1), round((y2 - y1) / ch * 100, 1))
            # 같은 글자가 화면 안에 여러 개 있으면 멀리 떨어진 것끼리 묶여 커진다.
            # 그런 상자는 무엇을 가리키는지 알 수 없다 — 글자를 더 정확히 적게 알린다.
            if 상[2] * 상[3] > 55 * 55:
                빠짐.append(f"{이름}/{별명} — 상자가 너무 큽니다 {상} · "
                            f"찾은 도형 {len(상자)}개. 가리킬 글자를 더 좁혀 적으세요")
                continue
            칸[별명] = 상
        결과[이름] = 칸

    if 빠짐:
        print("찾지 못한 대상:", file=sys.stderr)
        for m in 빠짐:
            print("  " + m, file=sys.stderr)
        return 1

    줄 = ['# -*- coding: utf-8 -*-',
          '"""노션 화면에서 가리킬 자리 — **계산해서 넣은 값이다. 손으로 고치지 않는다.**',
          '',
          '    python 제작/빌드도구/노션좌표뽑기.py',
          '',
          '원본 pptx 의 도형 위치에서 뽑았다. 가리킬 대상을 바꾸려면 그 도구의',
          '`대상` 목록을 고치고 다시 돌린다.',
          '"""',
          '',
          '자리 = {']
    for 이름, 칸 in 결과.items():
        줄.append(f'    "{이름}": {{')
        for 별명, v in 칸.items():
            줄.append(f'        "{별명}": {v},')
        줄.append("    },")
    줄.append("}")
    결과파일.write_text("\n".join(줄) + "\n", encoding="utf-8")

    print(f"  {sum(len(v) for v in 결과.values())}개 자리 → {결과파일.relative_to(REPO)}")
    for 이름, 칸 in 결과.items():
        print(f"  {이름}")
        for 별명, v in 칸.items():
            print(f"      {별명:<12} {v}")
    return 0


if __name__ == "__main__":
    sys.exit(main())

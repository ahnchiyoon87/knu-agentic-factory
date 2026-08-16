# -*- coding: utf-8 -*-
"""슬라이드를 **요소별로 편집되는** pptx 로 만든다.

    python 제작/빌드도구/편집pptx.py [노션|2일차|3일차]

왜 이렇게 만드나
    예전에는 슬라이드 한 장을 통째로 PNG 하나로 넣었다. 파워포인트에서 열어도
    **손댈 데가 없었다.** 강조 상자가 1mm 어긋나도 코드를 고쳐 다시 뽑아야 했다.

    이제 화면에 있는 것을 그대로 파워포인트 요소로 만든다.
      글자        → 텍스트 상자 (고쳐 쓸 수 있다)
      캡처        → 그림
      주황 상자   → 사각형 도형 (끌어서 옮기면 된다)
      번호 동그라미 → 타원 + 글자
      삽화(SVG)   → 그림 한 장

    좌표는 손으로 계산하지 않는다. **브라우저가 실제로 그린 위치를 읽어 온다.**
    (직접 계산하면 HTML 미리보기와 pptx 가 서서히 어긋난다. 실제로 그랬다.)

    1600×900 로 그린 화면을 13.333in×7.5in 슬라이드에 그대로 옮긴다.
    1px = 7620 EMU · 글자 1px = 0.6pt.
"""

from __future__ import annotations

import sys as _sys

for _s in (_sys.stdout, _sys.stderr):
    if (getattr(_s, "encoding", "") or "").lower().replace("-", "") != "utf8":
        try:
            _s.reconfigure(encoding="utf-8", errors="replace")
        except Exception:
            pass

import base64
import io
import json
import re
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parents[1]
덱 = REPO / "제작" / "산출물" / "덱"
강의 = REPO / "강의"

EMU = 7620          # 1px
PT = 0.6            # 1px → pt

묶음 = {
    "노션": ("노션덱.html", "2일차", "노션_자산화.pptx"),
    "2일차": ("2일차덱.html", "2일차", "슬라이드.pptx"),
    "3일차": ("3일차덱.html", "3일차", "슬라이드.pptx"),
}

# ── 브라우저가 그린 것을 그리기 명령으로 바꾼다 ────────────────────────────
# 요소를 훑으며 (1) 배경·테두리 (2) 글자 (3) 그림 (4) 삽화 를 뽑는다.
# 글자는 **더 이상 쪼갤 수 없는 칸**에서만 뽑는다 — 안에 또 칸이 있으면 내려간다.
훑기 = r"""
(idx) => {
  // 앞 장에 남은 표식을 지운다 — 안 지우면 숨은 삽화를 찍으려다 30초씩 기다린다
  document.querySelectorAll('svg[data-svgn]').forEach(s => delete s.dataset.svgn);
  const sections = [...document.querySelectorAll('section.s')];
  sections.forEach((s, i) => { s.style.display = i === idx ? 'flex' : 'none';
                               s.style.margin = '0'; });
  window.scrollTo(0, 0);
  const sec = sections[idx];
  const base = sec.getBoundingClientRect();
  const ops = [];
  let svgn = 0;

  const px = v => parseFloat(v) || 0;
  const rel = r => ({x: r.left - base.left, y: r.top - base.top,
                     w: r.width, h: r.height});
  const 보임 = c => c && c !== 'transparent' && !/rgba\(0,\s*0,\s*0,\s*0\)/.test(c);
  const 블록 = el => {
    const d = getComputedStyle(el).display;
    return /block|flex|grid|table|list-item/.test(d);
  };

  // 글자 조각 — 굵기·색·고정폭·밑칠이 바뀔 때마다 끊는다
  function runs(el) {
    const out = [];
    (function walk(node, st) {
      for (const ch of node.childNodes) {
        if (ch.nodeType === 3) {
          if (ch.textContent) out.push({t: ch.textContent, ...st});
        } else if (ch.nodeType === 1) {
          if (ch.tagName === 'BR') { out.push({t: '\n', ...st}); continue; }
          const c = getComputedStyle(ch);
          walk(ch, {
            b: parseInt(c.fontWeight) >= 600,
            m: /Consolas|monospace|Courier/i.test(c.fontFamily),
            c: c.color,
            sz: px(c.fontSize),
          });
        }
      }
    })(el, (() => { const c = getComputedStyle(el);
                    return {b: parseInt(c.fontWeight) >= 600,
                            m: /Consolas|monospace|Courier/i.test(c.fontFamily),
                            c: c.color, sz: px(c.fontSize)}; })());
    return out;
  }

  // 글자 뒤에 깔린 색(밑칠·코드칸) — 줄마다 네모를 따로 그린다
  function 밑칠(el) {
    const r = [];
    for (const sp of el.querySelectorAll('span, code, em')) {
      const c = getComputedStyle(sp);
      if (!보임(c.backgroundColor)) continue;
      for (const box of sp.getClientRects())
        r.push({...rel(box), fill: c.backgroundColor,
                round: px(c.borderRadius),
                line: 보임(c.borderTopColor) && px(c.borderTopWidth) > 0
                      ? c.borderTopColor : null});
    }
    return r;
  }

  (function walk(el) {
    const c = getComputedStyle(el);
    if (c.display === 'none' || c.visibility === 'hidden') return;
    const r = rel(el.getBoundingClientRect());
    if (r.w < 0.5 || r.h < 0.5) return;

    if (el.tagName === 'IMG') {
      ops.push({k: 'img', ...r, src: el.src});
      if (px(c.borderTopWidth) > 0 && 보임(c.borderTopColor))   // 캡처 테두리
        ops.push({k: 'frame', ...r, fill: c.borderTopColor,
                  lw: px(c.borderTopWidth), round: px(c.borderRadius)});
      return;
    }
    if (el.tagName === 'svg') { el.dataset.svgn = svgn;
                                ops.push({k: 'svg', ...r, n: svgn++}); return; }

    // 이 칸 자체의 배경과 테두리 — .bar, 표 머리, 구분선, 동그라미가 여기서 나온다
    // 슬라이드 판(sec) 자체의 바탕은 도형으로 만들지 않는다 — 배경색으로 넣는다
    if (보임(c.backgroundColor) && el !== document.body && el !== sec)
      ops.push({k: 'fill', ...r, fill: c.backgroundColor, round: px(c.borderRadius),
                oval: c.borderRadius.includes('50%')});
    // 네 면이 다 같은 테두리면 **선 네 개가 아니라 사각형 하나**로 만든다.
    // 선으로 쪼개 두면 강조 상자를 옮길 때마다 네 번 잡아야 한다.
    const 면들 = ['Top', 'Bottom', 'Left', 'Right'].map(면 => ({
      w: px(c['border' + 면 + 'Width']), col: c['border' + 면 + 'Color'],
      st: c['border' + 면 + 'Style']}));
    const 살아있는 = 면들.filter(f => f.w > 0 && 보임(f.col) && f.st !== 'none');
    const 사방 = 살아있는.length === 4 &&
                 살아있는.every(f => f.w === 살아있는[0].w && f.col === 살아있는[0].col);
    if (사방) {
      ops.push({k: 'frame', ...r, fill: 살아있는[0].col, lw: 살아있는[0].w,
                round: px(c.borderRadius),
                dash: 살아있는[0].st === 'dashed'});
    } else {
      for (const [i, [면, dx, dy, dw, dh]] of [['Top', 0, 0, r.w, 0],
                                               ['Bottom', 0, r.h, r.w, 0],
                                               ['Left', 0, 0, 0, r.h],
                                               ['Right', r.w, 0, 0, r.h]].entries()) {
        const f = 면들[i];
        if (f.w > 0 && 보임(f.col) && f.st !== 'none')
          ops.push({k: 'line', x: r.x + dx, y: r.y + dy, w: dw, h: dh,
                    fill: f.col, lw: f.w, dash: f.st === 'dashed'});
      }
    }

    const 자식 = [...el.children];
    if (자식.length && 자식.some(블록)) { 자식.forEach(walk); return; }

    const 글 = el.textContent.trim();
    if (글) {
      for (const b of 밑칠(el)) ops.push({k: 'fill', ...b});
      // 칸 안쪽 여백을 빼야 글자가 제자리에 온다 — 각주(padding-top 26)·표 칸(24/30)이
      // 이것 때문에 통째로 위로 밀려 밑칠과 어긋났다.
      const [pl, pt, pr, pb] = ['Left', 'Top', 'Right', 'Bottom']
                                 .map(s => px(c['padding' + s]));
      const 가운데 = /flex|grid/.test(c.display) && /center/.test(c.alignItems);
      ops.push({k: 'text', x: r.x + pl, y: r.y + pt,
                w: Math.max(r.w - pl - pr, 1), h: Math.max(r.h - pt - pb, 1),
                runs: runs(el),
                sz: px(c.fontSize), lh: px(c.lineHeight) || px(c.fontSize) * 1.2,
                align: /center/.test(c.justifyContent) ? 'center' : c.textAlign,
                가운데: 가운데});
    } else 자식.forEach(walk);
  })(sec);

  const bg = getComputedStyle(sec).backgroundColor;
  return {ops, bg: 보임(bg) ? bg : null, w: base.width, h: base.height};
}
"""


def _한상자로(ops):
    """세로로 이어지는 글 조각을 **하나의 텍스트 상자**로 묶는다.

    문단마다 상자를 따로 만들면 파워포인트에서 옮길 때 하나씩 따로 잡힌다.
    한 덩어리로 묶고 문단 사이는 **줄 간격**으로 벌린다.

    묶는 조건 — 왼쪽 끝과 폭이 같고, 바로 아래에 있고, 사이에 **다른 것이 없을 것.**
    (구분선이나 그림이 끼어 있으면 다른 덩어리다)
    """
    나온것, 묶음, 사이 = [], [], []

    def 내보내기():
        nonlocal 묶음
        if len(묶음) == 1:
            나온것.append(묶음[0])
        elif 묶음:
            나온것.append({**묶음[0], "묶음": 묶음})
        묶음 = []

    for o in ops:
        if o["k"] != "text":
            사이.append(o)          # 판정은 다음 글 조각이 올 때 한다
            continue
        앞 = 묶음[-1] if 묶음 else None
        # 밑칠(fill)만 끼어 있으면 같은 덩어리다. 구분선·그림이 있으면 다른 덩어리.
        깨끗 = all(x["k"] == "fill" for x in 사이)
        붙나 = (앞 is not None and 깨끗
                and not o.get("가운데") and not 앞.get("가운데")
                and abs(o["x"] - 앞["x"]) < 2 and abs(o["w"] - 앞["w"]) < 4
                and o.get("align") == 앞.get("align")
                and 0 <= o["y"] - (앞["y"] + 앞["h"]) < 붙는틈)
        if not 붙나:
            내보내기()
        나온것.extend(사이); 사이 = []
        묶음.append(o)
    내보내기()
    나온것.extend(사이)
    return 나온것


붙는틈 = 26.0      # 문단 사이가 이보다 벌어지면 다른 덩어리로 본다


def _색(s: str):
    from pptx.dml.color import RGBColor
    m = re.findall(r"[\d.]+", s or "")
    if len(m) >= 3:
        return RGBColor(int(float(m[0])), int(float(m[1])), int(float(m[2])))
    return RGBColor(0, 0, 0)


def _투명한가(s: str) -> bool:
    m = re.findall(r"[\d.]+", s or "")
    return len(m) >= 4 and float(m[3]) < 0.05


def 한장(prs, 명령, 바탕, 삽화: dict[int, bytes]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Pt

    s = prs.slides.add_slide(prs.slide_layouts[6])

    # 바탕은 **도형이 아니라 슬라이드 배경색**으로 넣는다.
    # 흰 사각형을 한 장 깔아 두면 요소를 누를 때마다 그 판이 먼저 잡히고
    # 크기 조절 손잡이도 판 것이 잡혀 아무것도 못 만진다. 실제로 그랬다.
    if 바탕:
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = _색(바탕)

    맞춤 = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY}

    for o in 명령:
        k = o["k"]
        X, Y = Emu(round(o["x"] * EMU)), Emu(round(o["y"] * EMU))
        W, H = Emu(round(max(o["w"], 1) * EMU)), Emu(round(max(o["h"], 1) * EMU))

        if k == "fill":
            if _투명한가(o["fill"]):
                continue
            모양 = (MSO_SHAPE.OVAL if o.get("oval") else
                    MSO_SHAPE.ROUNDED_RECTANGLE if o.get("round", 0) >= 3 else
                    MSO_SHAPE.RECTANGLE)
            sh = s.shapes.add_shape(모양, X, Y, W, H)
            sh.fill.solid(); sh.fill.fore_color.rgb = _색(o["fill"])
            if o.get("line"):
                sh.line.color.rgb = _색(o["line"]); sh.line.width = Pt(0.75)
            else:
                sh.line.fill.background()
            sh.shadow.inherit = False
            if 모양 == MSO_SHAPE.ROUNDED_RECTANGLE:      # 모서리를 CSS 만큼만 둥글게
                sh.adjustments[0] = min(0.5, o["round"] / max(o["w"], o["h"], 1))

        elif k == "line":
            from pptx.enum.shapes import MSO_CONNECTOR
            ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, X, Y,
                                        Emu(round((o["x"] + o["w"]) * EMU)),
                                        Emu(round((o["y"] + o["h"]) * EMU)))
            ln.line.color.rgb = _색(o["fill"])
            ln.line.width = Pt(max(o["lw"] * PT, 0.5))
            if o.get("dash"):
                from pptx.enum.dml import MSO_LINE_DASH_STYLE
                ln.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        elif k == "frame":
            sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if o["round"] >= 3
                                    else MSO_SHAPE.RECTANGLE, X, Y, W, H)
            sh.fill.background()
            sh.line.color.rgb = _색(o["fill"]); sh.line.width = Pt(o["lw"] * PT)
            sh.shadow.inherit = False
            if o.get("dash"):
                from pptx.enum.dml import MSO_LINE_DASH_STYLE
                sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            if o["round"] >= 3:
                sh.adjustments[0] = min(0.5, o["round"] / max(o["w"], o["h"], 1))

        elif k in ("img", "svg"):
            자료 = (삽화[o["n"]] if k == "svg"
                    else base64.b64decode(o["src"].split(",", 1)[1]))
            s.shapes.add_picture(io.BytesIO(자료), X, Y, W, H)

        elif k == "text":
            가운데 = o.get("가운데")
            덩어리 = o.get("묶음") or [o]
            끝 = 덩어리[-1]
            높이 = 끝["y"] + 끝["h"] - o["y"]
            # 글상자는 CSS 칸보다 조금 넓게 — 글꼴이 미세하게 달라 줄이 일찍 꺾이는 것을 막는다
            tb = s.shapes.add_textbox(X, Y if 가운데 else Emu(round((o["y"] - 2) * EMU)),
                                      Emu(round((o["w"] + (0 if 가운데 else 8)) * EMU)),
                                      Emu(round((높이 + (0 if 가운데 else 6)) * EMU)))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE if 가운데 else MSO_ANCHOR.TOP
            첫문단 = True
            p = tf.paragraphs[0]
            for 자리, 조각들 in enumerate(덩어리):
                if not 첫문단:
                    p = tf.add_paragraph()
                    # 문단 사이는 **빈 줄이 아니라 간격으로** 벌린다
                    앞 = 덩어리[자리 - 1]
                    p.space_before = Pt(max(조각들["y"] - (앞["y"] + 앞["h"]), 0) * PT)
                p.alignment = 맞춤.get(조각들["align"], PP_ALIGN.LEFT)
                p.line_spacing = Pt(조각들["lh"] * PT)
                첫문단 = False
                p = _글넣기(tf, p, 조각들, 맞춤)
    return s


def _글넣기(tf, p, o, 맞춤):
    """한 문단을 채운다. <br> 이 있으면 그 자리에서 문단을 하나 더 연다."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    for 조각 in o["runs"]:
        for i, 토막 in enumerate(조각["t"].split("\n")):
            if i:
                p = tf.add_paragraph()
                p.alignment = 맞춤.get(o["align"], PP_ALIGN.LEFT)
                p.line_spacing = Pt(o["lh"] * PT)
            if not 토막:
                continue
            r = p.add_run(); r.text = 토막
            f = r.font
            f.size = Pt(round(조각["sz"] * PT, 1))
            f.bold = bool(조각["b"])
            f.color.rgb = _색(조각["c"])
            f.name = "Consolas" if 조각["m"] else "맑은 고딕"
    return p


def main() -> int:
    이름 = sys.argv[1] if len(sys.argv) > 1 else "노션"
    if 이름 not in 묶음:
        print(f"모르는 덱: {이름} — {' · '.join(묶음)}", file=sys.stderr)
        return 1
    html, 일차, 파일 = 묶음[이름]
    원본 = 덱 / html
    if not 원본.is_file():
        print(f"덱이 없습니다: {원본} — 먼저 build_{이름}.py 를 돌리세요", file=sys.stderr)
        return 1

    from playwright.sync_api import sync_playwright
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(1600 * EMU)
    prs.slide_height = Emu(900 * EMU)

    with sync_playwright() as pw:
        b = pw.chromium.launch()
        pg = b.new_page(viewport={"width": 1600, "height": 900},
                        device_scale_factor=2)
        pg.goto(원본.as_uri())
        pg.wait_for_timeout(400)
        수 = pg.evaluate("() => document.querySelectorAll('section.s').length")

        for i in range(수):
            결과 = pg.evaluate(훑기, i)
            결과["ops"] = _한상자로(결과["ops"])
            # 삽화(SVG)는 도형으로 못 옮긴다 — 그 자리만 2배로 찍어 그림으로 넣는다
            삽화 = {}
            for el in pg.query_selector_all("svg[data-svgn]"):
                n = int(el.get_attribute("data-svgn"))
                try:
                    삽화[n] = el.screenshot(omit_background=True, timeout=5000)
                except Exception:
                    삽화[n] = pg.screenshot(clip=el.bounding_box())
            한장(prs, 결과["ops"], 결과["bg"], 삽화)
            print(f"    {i + 1:>2}/{수}", end="\r")
        b.close()

    dst = 강의 / 일차 / "강의자료" / 파일
    dst.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst))
    print(f"  만듦  {일차}/{dst.name}  ({수}장 · {dst.stat().st_size // 1024}KB)")
    print("  글자·주황 상자·번호가 전부 파워포인트 요소입니다 — 그대로 끌어 옮기면 됩니다.")
    return 0


if __name__ == "__main__":
    sys.exit(main())

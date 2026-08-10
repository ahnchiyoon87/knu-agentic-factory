# -*- coding: utf-8 -*-
"""노션 자산화 실습 덱 (2일차 1교시 · 60분) — PPTX 조립기

Day 1 덱과 같은 규격: 16:9, 맑은 고딕, 순백 배경, 표지형만 남색.
노션 화면 목업은 도형으로 직접 그린다 (NotebookLM 으로는 안 되는 부분).

    python build_notion_pptx.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 규격 (Day 1 덱과 동일) ────────────────────────────────────────────────
INK    = RGBColor(0x1B, 0x1F, 0x26)
INK2   = RGBColor(0x3B, 0x43, 0x50)
NAVY   = RGBColor(0x12, 0x30, 0x5A)
WASH   = RGBColor(0xFB, 0xE3, 0xD6)
TINT   = RGBColor(0xEE, 0xF3, 0xF9)
LINE   = RGBColor(0xE2, 0xE6, 0xEB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x1F, 0x6F, 0x62)
RED    = RGBColor(0xC0, 0x39, 0x2B)
AMBER  = RGBColor(0xC7, 0x7A, 0x28)
GREY   = RGBColor(0x9B, 0x9A, 0x97)
NBORD  = RGBColor(0xE7, 0xE6, 0xE3)   # 노션 테두리
NTEXT  = RGBColor(0x37, 0x35, 0x2F)   # 노션 글자
NBLUE  = RGBColor(0x23, 0x83, 0xE2)
TAG_R  = RGBColor(0xFF, 0xE2, 0xDD)
TAG_RT = RGBColor(0x5D, 0x17, 0x15)
TAG_G  = RGBColor(0xDB, 0xED, 0xDB)
TAG_GT = RGBColor(0x1C, 0x38, 0x29)
TAG_Y  = RGBColor(0xFA, 0xDE, 0xC9)
TAG_YT = RGBColor(0x49, 0x29, 0x0E)

SW, SH = Inches(13.333), Inches(7.5)
ML = Inches(0.8)
CW = Inches(11.733)
FONT = "맑은 고딕"
MONO = "Consolas"
T_TITLE, T_BODY, T_BIG, T_SMALL = 27, 16.5, 33, 12.5

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]
N = 0


def _font(run, size, bold=False, color=INK, font=FONT):
    f = run.font
    f.name = font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", font)


def tb(sl, x, y, w, h, text, size, bold=False, color=INK, align=PP_ALIGN.LEFT,
       anchor=MSO_ANCHOR.TOP, ls=1.15, sa=6, font=FONT):
    box = sl.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    for i, ln in enumerate(text if isinstance(text, list) else [text]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = ls; p.space_after = Pt(sa)
        r = p.add_run(); r.text = ln
        _font(r, size, bold, color, font)
    return box


def rect(sl, x, y, w, h, fill, line=None, radius=False):
    shp = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if radius:
        try: shp.adjustments[0] = 0.06
        except Exception: pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def hline(sl, x, y, w, color=LINE, weight=1.0):
    ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(weight))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background(); ln.shadow.inherit = False


def new(title=None, sub=None, tag=None):
    """일반 슬라이드. 머리말을 얹고 본문 시작 y 를 돌려준다."""
    global N
    N += 1
    sl = prs.slides.add_slide(BLANK)
    y = Inches(0.62)
    if tag:
        tb(sl, ML, Inches(0.42), CW, Inches(0.3), tag, T_SMALL, True, GREEN)
        y = Inches(0.78)
    if title:
        tb(sl, ML, y, CW, Inches(0.75), title, T_TITLE, True, INK)
        y += Inches(0.85)
    if sub:
        tb(sl, ML, y, CW, Inches(0.5), sub, T_BODY, False, INK2)
        y += Inches(0.55)
    tb(sl, Inches(12.5), Inches(6.95), Inches(0.6), Inches(0.3), str(N),
       T_SMALL, False, GREY, PP_ALIGN.RIGHT)
    return sl, y + Inches(0.1)


def cover(eyebrow, lines, sub):
    global N
    N += 1
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, SW, SH, NAVY)
    tb(sl, ML, Inches(2.3), CW, Inches(0.4), eyebrow, 13, True, RGBColor(0x9D, 0xB4, 0xD4))
    tb(sl, ML, Inches(2.8), CW, Inches(1.8), lines, 40, True, WHITE, ls=1.22)
    tb(sl, ML, Inches(4.9), CW, Inches(0.5), sub, T_BODY, False, RGBColor(0xC3, 0xD2, 0xE6))
    return sl


def bullets(sl, y, items, size=T_BODY, gap=0.46, color=INK, bullet="·"):
    for it in items:
        bold = it.startswith("**")
        txt = it.replace("**", "")
        tb(sl, ML, y, CW, Inches(0.42), f"{bullet} {txt}" if bullet else txt,
           size, bold, color)
        y += Inches(gap)
    return y


def big(sl, y, text, size=T_BIG, color=INK, align=PP_ALIGN.LEFT):
    tb(sl, ML, y, CW, Inches(1.2), text, size, True, color, align, ls=1.25)
    return y + Inches(0.55 * (len(text) if isinstance(text, list) else 1) + 0.4)


def quote(sl, y, text, color=GREEN):
    """왼쪽 굵은 세로선 + 큰 글씨."""
    h = Inches(0.55 * (len(text) if isinstance(text, list) else 1) + 0.35)
    rect(sl, ML, y, Inches(0.06), h, color)
    tb(sl, ML + Inches(0.28), y + Inches(0.05), CW - Inches(0.4), h, text, 24, True, INK, ls=1.3)
    return y + h + Inches(0.25)


def note(sl, y, text, fill=WASH, color=INK2, size=15.5):
    lines = text if isinstance(text, list) else [text]
    h = Inches(0.36 * len(lines) + 0.3)
    rect(sl, ML, y, CW, h, fill)
    tb(sl, ML + Inches(0.22), y + Inches(0.14), CW - Inches(0.44), h, lines, size, False, color)
    return y + h + Inches(0.22)


def table(sl, y, header, rows, widths, row_h=0.5):
    x = ML
    tot = sum(widths)
    ws = [CW * (w / tot) for w in widths]
    rect(sl, ML, y, CW, Inches(0.46), TINT)
    for i, htxt in enumerate(header):
        tb(sl, x + Inches(0.12), y + Inches(0.08), ws[i], Inches(0.3), htxt, 13.5, True, INK2)
        x += ws[i]
    y += Inches(0.46)
    for r in rows:
        x = ML
        hline(sl, ML, y, CW)
        for i, c in enumerate(r):
            bold = c.startswith("**")
            col = INK
            txt = c.replace("**", "")
            if txt.startswith("!R"): col, txt = RED, txt[2:]
            elif txt.startswith("!G"): col, txt = GREEN, txt[2:]
            tb(sl, x + Inches(0.12), y + Inches(0.09), ws[i], Inches(0.35), txt, 15, bold, col)
            x += ws[i]
        y += Inches(row_h)
    hline(sl, ML, y, CW)
    return y + Inches(0.25)


# ── 노션 목업 부품 ────────────────────────────────────────────────────────
def nframe(sl, x, y, w, h, url):
    """브라우저 프레임."""
    rect(sl, x, y, w, h, WHITE, NBORD)
    rect(sl, x, y, w, Inches(0.34), RGBColor(0xEA, 0xED, 0xF0))
    for i in range(3):
        d = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.14 + i * 0.2), y + Inches(0.1),
                                Inches(0.11), Inches(0.11))
        d.fill.solid(); d.fill.fore_color.rgb = RGBColor(0xCF, 0xD4, 0xDA)
        d.line.fill.background(); d.shadow.inherit = False
    tb(sl, x + Inches(0.85), y + Inches(0.05), w - Inches(1.0), Inches(0.26), url,
       11.5, False, GREY, font=MONO)
    return y + Inches(0.34)


def ntag(sl, x, y, text, kind="r"):
    fill, col = {"r": (TAG_R, TAG_RT), "g": (TAG_G, TAG_GT), "y": (TAG_Y, TAG_YT)}[kind]
    w = Inches(0.16 + 0.115 * len(text))
    rect(sl, x, y, w, Inches(0.26), fill, radius=True)
    tb(sl, x, y + Inches(0.01), w, Inches(0.24), text, 11.5, True, col, PP_ALIGN.CENTER)
    return w


def ndb(sl, x, y, w, title, header, rows, foot=None):
    """노션 인라인 데이터베이스 목업. rows = [(제목, 태그종류, 태그글, 자료)]"""
    tb(sl, x, y, w, Inches(0.32), title, 15, True, NTEXT)
    y += Inches(0.4)
    cols = [w * 0.56, w * 0.26, w * 0.18]
    cx = x
    for i, h in enumerate(header):
        tb(sl, cx, y, cols[i], Inches(0.28), h, 12, False, GREY)
        cx += cols[i]
    y += Inches(0.3)
    hline(sl, x, y, w, NBORD)
    for r in rows:
        y += Inches(0.06)
        tb(sl, x, y, cols[0], Inches(0.3), r[0], 13.5, True, NTEXT)
        ntag(sl, x + cols[0], y + Inches(0.02), r[2], r[1])
        tb(sl, x + cols[0] + cols[1], y, cols[2], Inches(0.3), r[3], 13, False, GREY)
        y += Inches(0.36)
        hline(sl, x, y, w, NBORD)
    if foot:
        tb(sl, x, y + Inches(0.1), w, Inches(0.3), foot, 12.5, False, GREY)
        y += Inches(0.36)
    return y


def ncode(sl, x, y, w, lines, lang="Python"):
    h = Inches(0.26 * len(lines) + 0.42)
    rect(sl, x, y, w, h, RGBColor(0xF7, 0xF6, 0xF3), NBORD)
    tb(sl, x + w - Inches(1.0), y + Inches(0.06), Inches(0.9), Inches(0.24), lang + " ▾",
       10.5, False, GREY, PP_ALIGN.RIGHT)
    tb(sl, x + Inches(0.16), y + Inches(0.2), w - Inches(0.3), h, lines, 12, False,
       NTEXT, ls=1.35, sa=0, font=MONO)
    return y + h + Inches(0.12)


# ══════════════════════════════════════════════════════════════════════════
# 슬라이드 (샘플 4장 — 틀 확인용)
# ══════════════════════════════════════════════════════════════════════════

# S1 표지
cover("2일차 · 1교시", ["오늘 배운 걸", "내 것으로 만들기"], "노션으로 자산화하기 · 60분")

# S2 큰 문장형
sl, y = new(tag="왜 하는가")
y = big(sl, y + Inches(0.5), ["3주 뒤에 오늘 배운 게", "얼마나 남아 있을까요?"], 38)
note(sl, y + Inches(0.3), "실험 결과는 잔인합니다.", TINT, INK2, 17)

# S3 목업형
sl, y = new("지금 여러분 화면", "카드 제목만 있고, 전부 빨간색입니다", tag="구멍 뽑기")
by = nframe(sl, ML, y, CW, Inches(3.5), "notion.so — Agentic AI 특강 · 2일차")
ndb(sl, ML + Inches(0.45), by + Inches(0.25), CW - Inches(0.9), "💡 개념 카드",
    ["Aa 개념", "◐ 이해도", "≡ 자료"],
    [("윈도 크기는 어떻게 정하나", "r", "설명못함", "2p"),
     ("왜 지금 값을 평균에서 빼는가", "r", "설명못함", "2p"),
     ("오탐과 미탐은 무엇이 다른가", "r", "설명못함", "3p"),
     ("결측 처리에 정답이 없는 이유", "r", "설명못함", "4p")],
    "＋ 새로 만들기")
note(sl, Inches(6.35), "전부 빨간색이죠. 부끄러운 게 아니라 정확한 겁니다.", WASH, INK, 17)

# S4 표형
sl, y = new("물음표를 카드 제목으로", tag="구멍 뽑기")
table(sl, y, ["필기에 있던 것", "카드 제목"],
      [["\"W를 몇으로 ????\"", "**윈도 크기는 어떻게 정하나"],
       ["\"지금값은 평균에 넣지말라 — 무슨말?\"", "**왜 지금 값을 평균에서 빼는가"],
       ["\"오탐 미탐 헷갈림\"", "**오탐과 미탐은 무엇이 다른가"],
       ["\"뭘 골라야 하는거임?\"", "**결측 처리에 정답이 없는 이유"]],
      [52, 48])
note(sl, Inches(5.3), "지금은 제목만. 내용은 나중에 채웁니다.", TINT, INK2, 17)

out = Path(__file__).parent / "샘플_틀확인.pptx"
prs.save(out)
print(f"저장: {out} · {N}장")

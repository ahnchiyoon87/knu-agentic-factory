# -*- coding: utf-8 -*-
"""
Day 1 강의 1 (38장) PPTX 조립기

원칙
  - 글자는 설명란 원문을 그대로 옮긴 아래 SLIDES 데이터에서 찍는다 (NotebookLM 글자는 버림)
  - 삽화는 00_작업파일/삽화/art_NN.png (선정본에서 잘라낸 것)
  - 배경은 전부 순백, 표지형(1p·38p)만 남색 #12305A
  - 크기: 제목 28pt / 설명문·본문·결론 17pt / 큰 낱말 34pt / 라벨·용어 12.5pt
"""
import pathlib, sys, io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

ART = pathlib.Path("D:/work/study/경남대특강/작업장/슬라이드작업/00_작업파일/삽화")
CAP = pathlib.Path("D:/work/study/경남대특강/산출물/W1_팩토리시뮬레이터/docs/캡처")
OUT = pathlib.Path("D:/work/study/경남대특강/산출물/W2_강의슬라이드/Day1_강의1_38장.pptx")

NAVY   = RGBColor(0x12, 0x30, 0x5A)
NAVY2  = RGBColor(0x1E, 0x4E, 0x8C)
ORANGE = RGBColor(0xD4, 0x54, 0x1E)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
INK    = RGBColor(0x11, 0x11, 0x11)
INK2   = RGBColor(0x3B, 0x43, 0x50)
WASH   = RGBColor(0xFB, 0xE3, 0xD6)
TINT   = RGBColor(0xEE, 0xF3, 0xF9)
LINE   = RGBColor(0xE2, 0xE6, 0xEB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LBLUE  = RGBColor(0xC3, 0xD2, 0xE6)
LBLUE2 = RGBColor(0x9D, 0xB4, 0xD4)

SW, SH = Inches(13.333), Inches(7.5)
ML, MR = Inches(0.8), Inches(12.533)
CW = Inches(11.733)
FONT = "맑은 고딕"

T_TITLE, T_SUB, T_BODY, T_FOOT = 27, 16.5, 16.5, 16.5
T_WORD, T_WSUB, T_TAG, T_SMALL = 33, 15, 13, 12

def _set_font(run, size, bold=False, color=INK):
    f = run.font
    f.name = FONT; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    ea = run._r.rPr.get_or_change_to_latin() if False else None
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", FONT)

def tb(slide, x, y, w, h, text, size, bold=False, color=INK, align=PP_ALIGN.LEFT,
       fill=None, anchor=MSO_ANCHOR.TOP, line_spacing=1.12, space_after=6):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text if isinstance(text, list) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        r = p.add_run(); r.text = ln
        _set_font(r, size, bold, color)
    if fill is not None:
        box.fill.solid(); box.fill.fore_color.rgb = fill
        box.line.fill.background()
    return box

def rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(0.75)
    else: s.line.fill.background()
    s.shadow.inherit = False
    return s

def hline(slide, x, y, w, color=LINE, weight=1.0):
    ln = slide.shapes.add_connector(1, x, y, x + w, y)
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln

def pic_fit(slide, path, x, y, w, h, align="c"):
    im = Image.open(path); iw, ih = im.size
    r = min(w / iw, h / ih)
    pw, ph = int(iw * r), int(ih * r)
    px = x + (int(w) - pw) // 2 if align == "c" else x
    py = y + (int(h) - ph) // 2
    slide.shapes.add_picture(str(path), int(px), int(py), pw, ph)

def head(slide, title, sub, num=None, title_w=CW):
    y = Inches(0.5)
    if num:
        tb(slide, ML, Inches(0.42), Inches(0.8), Inches(0.32), num, 13, True, RGBColor(0xB9,0xC4,0xD2))
        y = Inches(0.72)
    tlines = title if isinstance(title, list) else [title]
    tb(slide, ML, y, title_w, Inches(0.55*len(tlines)), tlines, T_TITLE, True, INK, space_after=2)
    sy = y + Inches(0.58*len(tlines))
    slines = sub if isinstance(sub, list) else [sub]
    tb(slide, ML, sy, CW, Inches(0.34*len(slines)+0.1), slines, T_SUB, False, INK2, space_after=2)
    return sy + Inches(0.36*len(slines) + 0.12)

def foot(slide, lines):
    """lines: [(text, hl)] — 아래에서 위로 쌓는다. 긴 줄은 두 줄 높이. 반환: 규칙선 y"""
    hs = [0.66 if len(t_) > 55 else 0.38 for t_, _ in lines]
    top = 7.5 - 0.22 - sum(hs)
    hline(slide, ML, Inches(top - 0.06), CW)
    cy = top
    for (txt, hl), h in zip(lines, hs):
        tb(slide, ML, Inches(cy), CW, Inches(h), txt, T_FOOT, bool(hl), INK,
           fill=WASH if hl else None, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05, space_after=0)
        cy += h
    return Inches(top - 0.18)

def body_paras(slide, x, w, y, paras, size=T_BODY, gap=0.14):
    """문단 전부를 한 텍스트상자에 담는다 — 워드랩이 높이를 정하므로 겹침이 없다."""
    box = tb(slide, x, y, w, Inches(4.6), paras, size, False, RGBColor(0x22,0x22,0x22),
             line_spacing=1.18, space_after=int(gap*72))
    return y

# ─────────────────────────────────────────────────────────────
def cover(slide, eyebrow, title_lines, sub, tails):
    rect(slide, 0, 0, SW, SH, NAVY)
    tb(slide, ML, Inches(1.15), CW, Inches(0.4), eyebrow, 15, False, LBLUE2)
    tb(slide, ML, Inches(1.7), CW, Inches(0.75*len(title_lines)), title_lines, 36, True, WHITE, space_after=2)
    bar_y = 1.7 + 0.78 * len(title_lines) + 0.15
    rect(slide, ML, Inches(bar_y), Inches(1.3), Inches(0.09), ORANGE)
    tb(slide, ML, Inches(bar_y + 0.3), CW, Inches(0.4), sub, 16, False, LBLUE)
    ty = 7.5 - 0.55 - 0.42 * len(tails)
    for i, (txt, bold) in enumerate(tails):
        tb(slide, ML, Inches(ty + 0.42 * i), CW, Inches(0.4), txt, 15,
           bold, WHITE if bold else LBLUE2)
    for r_in in (1.0, 1.65, 2.3, 2.95, 3.5):
        d = Inches(r_in * 2)
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.2) - Inches(r_in), Inches(3.75) - Inches(r_in), d, d)
        o.fill.background(); o.line.color.rgb = RGBColor(0x3C, 0x55, 0x80); o.line.width = Pt(1.2)
        o.shadow.inherit = False

def rail_rows(slide, y0, y1, rows, icon_file=None, x_text=None, head_size=17, desc_size=15.5,
              head_colors=None, x0=ML, x1=MR):
    n = len(rows)
    rh = (y1 - y0) / n
    icon_w = Inches(1.15)
    if icon_file:
        pic_fit(slide, ART / icon_file, int(x0), int(y0), int(icon_w), int(y1 - y0))
        tx = x0 + icon_w + Inches(0.25)
    else:
        tx = x0
    tx = x_text if x_text else tx
    tw = x1 - tx
    for i, (h, d) in enumerate(rows):
        ry = y0 + int(rh * i)
        col = (head_colors[i] if head_colors else NAVY)
        tb(slide, tx, ry + Inches(0.02), tw, Inches(0.34), h, head_size, True, col)
        tb(slide, tx, ry + Inches(0.36), tw, int(rh) - Inches(0.36), d, desc_size, False, GRAY,
           line_spacing=1.08, space_after=0)
        if i < n - 1:
            hline(slide, tx, ry + int(rh) - Inches(0.02), tw)

def table_draw(slide, y0, header, rows, widths, hl_row=None, row_h=0.52, header_h=0.5,
               key_col=0, size=16):
    xs = [ML]
    for w in widths: xs.append(xs[-1] + Inches(w))
    rect(slide, ML, y0, CW, Inches(header_h), NAVY)
    for c, htxt in enumerate(header):
        tb(slide, xs[c] + Inches(0.1), y0 + Inches(0.07), Inches(widths[c]) - Inches(0.15),
           Inches(header_h - 0.1), htxt, size - 1, True, WHITE)
    cy = y0 + Inches(header_h)
    for ri, row in enumerate(rows):
        if hl_row is not None and ri == hl_row:
            rect(slide, ML, cy, CW, Inches(row_h), TINT)
        for c, cell in enumerate(row):
            bold = (c == key_col) or (hl_row is not None and ri == hl_row)
            col = NAVY if (hl_row is not None and ri == hl_row) else (INK if c == key_col else INK2)
            tb(slide, xs[c] + Inches(0.1), cy + Inches(0.05), Inches(widths[c]) - Inches(0.15),
               Inches(row_h - 0.08), cell, size, bold, col, line_spacing=1.0, space_after=0)
        cy += Inches(row_h)
        hline(slide, ML, cy, CW)
    return cy

# ═════════════════════ 38장 ═════════════════════
prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH
BLANK = prs.slide_layouts[6]

def new(): return prs.slides.add_slide(BLANK)

# 1
s = new()
cover(s, "DAY 1 · 보기", ["공장을 보는", "눈을 만든다"],
      "경남대 RISE 피지컬AI 사관학교 · 8월 Agentic AI 특강",
      [("피지컬 AI란 무엇인가 · 왜 하필 지금인가 · 에이전트란 무엇인가", False),
       ("오늘, 코드는 한 줄도 치지 않습니다.", True)])

# 2
s = new()
by = head(s, "새벽 3시, 아무도 없는 공장에서", "창원 K-정밀. 밤에도 공장은 돌고, 사람은 없다.")
fb = foot(s, [("사람이 옆에 서 있었더라도 몰랐을 속도다.", False)])
pic_fit(s, ART/"art_02.png", int(ML), int(by), int(Inches(5.3)), int(fb - by))
body_paras(s, Inches(6.5), Inches(6.0), by + Inches(0.5),
    ["새벽 3시, EQ-03이라는 CNC 한 대의 온도가 조금씩 오르기 시작했다.",
     "CNC는 컴퓨터 수치제어 공작기계. 창원 공장에 가장 흔하게 깔린 장비다.",
     "아침 7시, EQ-03은 여전히 돌고 있었다. 멈추지도, 불이 나지도, 경고음이 울리지도 않았다."])

# 3
s = new()
by = head(s, "문제는 그날 오후에 드러났습니다", "밤새 아무도 몰랐고, 아무 일도 없어 보였습니다.")
fb = foot(s, [('오늘 물을 것은 "왜 온도가 올랐나"가 아니다. 왜 네 시간 동안 아무도 몰랐나 — 이것이다.', True)])
body_paras(s, ML, Inches(5.4), by + Inches(0.4),
    ["밤새 깎아 낸 부품의 치수가 미세하게 틀어져 있었다. 하루치 생산분을 통째로 버렸다.",
     "설비가 고장 난 것도, 알람이 울린 것도 아니다. 그냥 조용히 나빠졌다."])
pic_fit(s, ART/"art_03.png", int(Inches(6.6)), int(by), int(Inches(5.9)), int(fb - by))

# 4
s = new()
by = head(s, "센서가 없어서 생긴 사고가 아닙니다", "EQ-03에는 온도계가 달려 있었고, 값은 계속 기록되고 있었다.")
fb = foot(s, [("이유는 셋입니다 — ① 볼 시간이 없다 ② 봐도 알아채지 못한다 ③ 볼 것이 너무 많다", False)])
pic_fit(s, ART/"art_04.png", int(ML), int(by), int(Inches(4.6)), int(fb - by))
tb(s, Inches(5.8), by + Inches(1.0), Inches(6.6), Inches(0.6), "볼 사람이 없었을 뿐이다", T_WORD, True, NAVY)
tb(s, Inches(5.8), by + Inches(1.75), Inches(6.6), Inches(0.4), "데이터는 1초마다 꼬박꼬박 쌓이고 있었다", T_SUB, False, GRAY)

# 5
s = new()
by = head(s, "사람은 볼 시간이 없습니다", "하루 스물네 시간 가운데, 사람이 지키고 있는 시간은 얼마나 될까요.", num="01")
fb = foot(s, [("사고는 사람이 가장 적은 시간에 납니다.", False)])
pic_fit(s, ART/"art_05.png", int(ML), int(by), int(Inches(5.3)), int(fb - by))
body_paras(s, Inches(6.5), Inches(6.0), by + Inches(0.6),
    ["K-정밀은 주간 근무만 하고, 야간에는 무인으로 돌아간다. 새벽 3시가 정확히 그 안에 있다.",
     "3교대로 사람을 붙여도 야간에는 인원이 줄고 집중도가 떨어진다. 사람을 더 뽑는 것으로는 안 풀린다."])

# 6
s = new()
by = head(s, "봐도 알아채지 못합니다", '62.1이 62.6으로 바뀌는 것을 보고 "이상하다"고 말할 사람은 없다.', num="02")
fb = foot(s, [("급격한 변화라면 누구나 안다. 문제는 이 사건이 급격하지 않았다는 것이다.", False),
              ("화면을 계속 보고 있었더라도 결과는 같았을 것이다.", False)])
tb(s, Inches(1.3), Inches(2.6), Inches(4.2), Inches(1.1), "62.1", 60, True, NAVY, align=PP_ALIGN.CENTER)
tb(s, Inches(6.1), Inches(2.7), Inches(1.1), Inches(0.9), "?", 40, True, ORANGE, align=PP_ALIGN.CENTER)
tb(s, Inches(7.9), Inches(2.6), Inches(4.2), Inches(1.1), "62.6", 60, True, NAVY, align=PP_ALIGN.CENTER)
tb(s, Inches(2.5), Inches(3.85), Inches(1.8), Inches(0.4), "새벽 3시", T_TAG, False, GRAY, align=PP_ALIGN.CENTER, fill=RGBColor(0xEF,0xF1,0xF3))
tb(s, Inches(9.1), Inches(3.85), Inches(1.8), Inches(0.4), "새벽 4시", T_TAG, False, GRAY, align=PP_ALIGN.CENTER, fill=RGBColor(0xEF,0xF1,0xF3))

# 7
s = new()
by = head(s, "볼 것이 너무 많습니다", "한 대만 보면 되는 게 아닙니다. 설비 6대 × 값 4가지 — 숫자 24개가 1초마다 바뀝니다.", num="03")
fb = foot(s, [("한 사람이 감당할 수 있는 양이 아니다.", False)])
pic_fit(s, ART/"art_07.png", int(ML), int(by), int(Inches(5.6)), int(fb - by))
body_paras(s, Inches(6.8), Inches(5.7), by + Inches(0.7),
    ["여기에 자율이동로봇 2대의 값까지 더해진다.",
     "여섯 대를 동시에, 쉬지 않고 봐야 한다. 사람 눈이 따라갈 양이 아니다."])

# 8
s = new()
by = head(s, "여러분이라면 어떻게 하시겠습니까?", '아마 가장 먼저 떠오르는 답은 이것일 겁니다 — "온도가 80도를 넘으면 경고음을 울리자."')
fb = foot(s, [("앞의 세 문제를 한 번에 푸는 것처럼 보입니다.", False)])
pic_fit(s, ART/"art_08.png", int(ML), int(by), int(Inches(4.8)), int(fb - by))
body_paras(s, Inches(6.0), Inches(6.5), by + Inches(0.7),
    ["공장에서 수십 년째 쓰는 방법이고, 실제로 쓸모도 있다. 갑자기 과열되는 상황은 이걸로 잡힌다.",
     "잠도 안 자고, 0.5도도 정확히 재고, 여섯 대를 동시에 본다."])

# 9
s = new()
by = head(s, "그런데 80도를 넘은 적이 없습니다", "선을 넘지 않고도 부품은 망가졌습니다.")
fb = foot(s, [('"값이 얼마인가"만 봐서는 부족하다. "값이 어떻게 변해 왔는가"를 봐야 한다.', True)])
body_paras(s, ML, Inches(5.2), by + Inches(0.25),
    ["그럼 선을 64도로 낮추면 되지 않나. 그러면 여름철에 정상적으로 돌던 다른 설비들이 하루 종일 울린다.",
     "경고가 자주 울리면 사람은 그것을 꺼 버린다. 그러면 처음 상태로 돌아간다.",
     "선을 높이면 못 잡고, 낮추면 다 울린다. 선 하나로는 빠져나갈 길이 없다."])
pic_fit(s, ART/"art_09.png", int(Inches(6.3)), int(by), int(Inches(6.2)), int(fb - by))

# 10
s = new()
by = head(s, "값이 아니라, 값이 걸어온 길",
          ["62.1 → 62.6 → 63.1 … 각각은 전부 정상입니다.",
           "그런데 네 시간 동안 한 번도 내려가지 않은 것은 정상이 아닙니다."])
fb = foot(s, [("셋 다 있어야 한다. 하나라도 빠지면 새벽 3시가 되풀이된다.", False)])
pic_fit(s, ART/"art_10.png", int(ML), int(by + Inches(0.05)), int(CW), int(fb - by - Inches(0.1)))

# 11
s = new()
by = head(s, "그리고 결정적인 네 번째", "행동한 결과를 다시 본다 — 그래야 고리가 닫힌다.")
fb = foot(s, [("결과를 안 보고 시킨 대로만 하는 구조는 고리가 열려 있다고 해서 개루프라 한다", False)])
pic_fit(s, ART/"art_11.png", int(ML), int(by), int(Inches(4.9)), int(fb - by))
tb(s, Inches(6.1), by + Inches(0.35), Inches(6.4), Inches(0.6), "폐루프", T_WORD, True, NAVY)
tb(s, Inches(6.1), by + Inches(1.02), Inches(6.4), Inches(0.35), "닫힌 고리 · closed loop", T_WSUB, False, GRAY)
body_paras(s, Inches(6.1), Inches(6.4), by + Inches(1.6),
    ["인지 → 판단 → 행동까지는 직선이다. 여기서 끝내면 시킨 것을 한 번 하고 마는 셈이다.",
     "속도를 낮췄으면 물어야 한다 — 온도가 정말 떨어졌는가? 안 떨어졌으면 더 낮추고, 떨어졌으면 그만한다."])

# 12
s = new()
by = head(s, "이것이 피지컬 AI입니다", "물리 세계를 센서로 읽고, 판단하고, 실제로 무언가를 움직이고, 그 결과를 다시 읽는 AI.")
fb = foot(s, [("읽는 것과 움직이는 것만 다를 뿐, 고리는 같습니다. 4일 뒤 이 고리를 돌리는 것은 여러분이 만든 프로그램입니다.", True)])
pic_fit(s, ART/"art_12.png", int(ML), int(by + Inches(0.05)), int(CW), int(fb - by - Inches(0.1)))

# 13
s = new()
by = head(s, "집에도 있습니다 — 에어컨과 전자레인지", "고리가 닫힌다는 게 무슨 뜻인지, 집에 있는 물건 두 개로 확인합니다.")
fb = foot(s, [('다시 공장으로 — 설비에 "속도를 20% 낮춰"라고 한 번 시키고 끝내면 그건 전자레인지다.', False),
              ("우리가 4일 동안 만들 것은 에어컨 쪽이다.", True)])
pic_fit(s, ART/"art_13.png", int(ML), int(by + Inches(0.05)), int(CW), int(fb - by - Inches(0.1)))

# 14
s = new()
by = head(s, "그런데 이 고리는 50년 전에도 있었습니다", "공장 자동화 장치는 1960년대부터 센서로 값을 읽고, 조건을 따지고, 밸브를 열고 닫았습니다.")
fb = foot(s, [("그렇게 오래된 구조인데, 왜 아직도 새벽 3시의 EQ-03을 아무도 몰랐나", True)])
pic_fit(s, ART/"art_14.png", int(ML), int(by), int(Inches(4.9)), int(fb - by))
tb(s, Inches(6.1), by + Inches(1.1), Inches(6.4), Inches(0.7), "50년", 44, True, NAVY)
tb(s, Inches(6.1), by + Inches(1.95), Inches(6.4), Inches(0.4), "이 구조가 존재해 온 기간", T_WSUB, False, GRAY)

# 15
s = new()
by = head(s, "막혀 있던 것은 가운데 「판단」이었습니다", "온도를 재는 일과 모터를 돌리는 일은 기계가 잘한다. 문제는 그 사이다.")
fb = foot(s, [("그래서 규칙을 계속 늘리게 된다. 100개, 500개. 서로 부딪히다 어느 순간 아무도 전체를 이해하지 못하는 덩어리가 된다.", False),
              ("50년 동안 이 칸에 넣을 수 있는 재료가 둘뿐이었다 — 이것이 병목의 정체다.", True)])
pic_fit(s, ART/"art_15.png", int(ML), int(by + Inches(0.05)), int(CW), int(fb - by - Inches(0.1)))

# 16
s = new()
by = head(s, "그 칸에 넣을 세 번째 재료가 생겼습니다",
          ["LLM — Large Language Model, 대규모 언어 모델.",
           "엄청난 양의 글을 읽고 다음에 올 말을 알아맞히도록 훈련된 프로그램."])
fb = foot(s, [('이것이 "왜 하필 지금인가"의 답이다. 센서가 좋아져서도, 로봇이 싸져서도 아니다. 50년 동안 막혀 있던 칸이 뚫렸기 때문이다.', True)])
pic_fit(s, ART/"art_16.png", int(ML), int(by), int(Inches(3.4)), int(fb - by))
tb(s, Inches(4.7), by + Inches(0.1), Inches(7.8), Inches(0.6), "LLM", T_WORD, True, NAVY)
tb(s, Inches(4.7), by + Inches(0.75), Inches(7.8), Inches(0.35), "범용 판단 엔진 · 특정 분야에 묶여 있지 않다", T_WSUB, False, GRAY)
body_paras(s, Inches(4.7), Inches(7.8), by + Inches(1.3),
    ["글을 다루려고 만들었는데, 훈련 과정에서 맥락을 보고 판단을 내리는 능력이 함께 생겼다.",
     '규칙을 500개 적어 두지 않아도 "이 온도 흐름이 평소와 다른가"를 맥락으로 본다.',
     "미리 예상하지 못한 상황에도 답을 낸다. 규칙은 침묵하지만 이쪽은 말을 한다.",
     '그리고 판단한 이유를 사람의 말로 설명한다. 규칙은 "조건 37번 걸림"이라고만 한다.'], gap=0.08)

# 17
s = new()
by = head(s, "다만, AI는 판단만 합니다", '"AI가 다 한다"는 말은 틀렸다. AI는 고리의 가운데 한 칸을 맡는다.')
fb = foot(s, [("4일 내내 이 경계를 지킨다. AI가 로그를 전부 읽거나 스스로 설비를 세운다고 말하지 않는다.", False)])
table_draw(s, by + Inches(0.15),
    ["누가", "무엇을 하는가"],
    [["센서", "온도·진동을 잰다"],
     ["프로그램(코드)", "값을 모으고 형식을 맞추고 저장한다"],
     ["제어 장치", "설비 속도를 실제로 바꾸고 로봇을 움직인다"],
     ["AI (LLM)", '모인 값을 보고 "이건 이상하다, 이렇게 하자"고 정하는 일만 한다']],
    [2.7, 9.033], hl_row=3, row_h=0.78, header_h=0.55)

# 18
s = new()
by = head(s, "기업이 찾는 사람은 AI 전문가가 아닙니다", "로봇도 센서도 공장에는 이미 깔려 있습니다. 장비가 없어서 못 하는 것이 아닙니다.")
fb = foot(s, [("오늘 만드는 관제 대시보드는 그 기업들이 요청한 과제 목록의 첫 줄과 같은 종류의 일이다. 수요조사 참여기업 15개사, 채용약정 4개사.", True)])
pic_fit(s, ART/"art_18.png", int(ML), int(by), int(CW), int(Inches(2.0)))
body_paras(s, ML, CW, by + Inches(2.15),
    ["기업이 실제로 호소하는 것은 이것이다 — 새 AI를 공장에 이미 깔려 있는 낡은 시스템과 이어 붙일 사람이 없다.",
     "그 다리를 놓을 수 있는 사람 — 이것이 참여기업이 원하는 인재상이다. 모델을 새로 만드는 연구자가 아니다."], gap=0.06)
tb(s, ML, by + Inches(3.15), CW, Inches(1.05),
   ["MES — Manufacturing Execution System, 제조실행시스템. 어떤 설비가 무엇을 얼마나 만들었는지 관리하는 프로그램",
    "PLC — Programmable Logic Controller. 설비를 직접 켜고 끄는 제어 장치",
    "SOP — Standard Operating Procedure. 표준 작업 절차서"], T_SMALL, False, GRAY, space_after=2)

# 19
s = new()
by = head(s, "그래서 4일 동안 이 고리를 한 토막씩 붙입니다", "오늘은 첫 토막 — 보는 눈을 만든다.")
fb = foot(s, [("넷째 날, 여러분이 만든 프로그램이 설비에 진짜로 명령을 내리면서 고리가 닫힌다.", False),
              ("그런데 그 판단을 LLM에게 그냥 물어보면 될까요?", True)])
table_draw(s, by + Inches(0.1),
    ["일자", "고리의 어느 부분", "만드는 것", "남는 것"],
    [["Day 1", "인지 — 보기", "설비·로봇 실시간 관제 대시보드", "관제화면"],
     ["Day 2", "인지의 구조화", "스펙 기반 관제 시스템, 알람 규칙과 작업지시", "DB 백엔드"],
     ["Day 3", "판단 — 알기", "센서 이상감지 + 진단 에이전트", "이상감지·진단"],
     ["Day 4", "행동 — 움직이기", "설비 제어와 로봇 파견 폐루프", "완성된 폐루프"]],
    [1.35, 2.5, 5.6, 2.283], hl_row=0, row_h=0.72, header_h=0.52)

# 20
s = new()
by = head(s, "그냥 물어보면 이렇게 됩니다", '"EQ-03 지금 온도 어때?" — 답은 매끄럽고, 틀렸다.')
fb = foot(s, [("아무리 판단을 잘해도 볼 수 없으면 소용이 없다. LLM에게는 눈이 없다.", True)])
pic_fit(s, ART/"art_20.png", int(ML), int(by), int(Inches(5.3)), int(fb - by))
body_paras(s, Inches(6.5), Inches(6.0), by + Inches(0.5),
    ["LLM은 훈련이 끝난 시점까지의 글을 읽었을 뿐이다. 창원 공장에 오늘 새벽 무슨 일이 있었는지는 그 글에 없다.",
     "K-정밀도, EQ-03도 LLM은 모른다.",
     "더 나쁜 것은 모른다고 말하지 않는다는 점이다. 그럴듯한 숫자를 지어내 문장을 완성한다."])

# 21
s = new()
by = head(s, "눈을 달아 줘도, 손이 없습니다", "그럼 온도 값을 복사해 붙여 넣으면 되지 않나 — 한 번은 됩니다.")
fb = foot(s, [("판단은 하는데 볼 눈도, 쓸 손도, 계속 돌 심장도 없다.", True)])
rail_rows(s, by + Inches(0.15), fb - Inches(0.1),
    [('"그럼 속도를 낮춰."', "문장만 만든다. 설비는 그대로 돈다."),
     ('"지난 일주일 이력도 봐."', "매번 사람이 찾아서 붙여 넣어야 한다."),
     ('"이상하면 알려 줘."', "물어볼 때만 대답한다. 스스로 계속 보지 않는다.")],
    icon_file="art_21.png")

# 22
s = new()
by = head(s, "그래서 에이전트입니다", "에이전트(Agent) — 대리인. 사람을 대신해 일을 맡아서 끝내는 존재.")
fb = foot(s, [("앞에서 그린 폐루프와 모양이 같다. 에이전트는 폐루프를 소프트웨어로 만든 것이다.", True)])
pic_fit(s, ART/"art_22.png", int(ML), int(by + Inches(0.05)), int(CW), int(fb - by - Inches(0.1)))

# 23
s = new()
by = head(s, "챗봇도 코파일럿도 아닙니다", "차이는 하나뿐이다 — 다음에 무엇을 할지 누가 정하는가.")
fb = foot(s, [("핵심 차이는 하나다. 사람이 정하면 챗봇·코파일럿, 자기가 정하면 에이전트다.", False)])
table_draw(s, by + Inches(0.1),
    ["구분", "챗봇", "코파일럿", "에이전트"],
    [["하는 일", "묻는 말에 답한다", "사람 옆에서 거든다", "일을 맡아서 끝낸다"],
     ["도구 사용", "안 쓴다", "사람이 시킬 때만", "스스로 골라 쓴다"],
     ["작동 횟수", "한 번", "한 번씩 여러 번", "끝날 때까지 반복"],
     ["다음을 누가 정하나", "사람", "사람", "자기가 정한다"],
     ["공장에 비유하면", "안내원", "보조원", "야간 당직자"]],
    [3.0, 2.7, 2.9, 3.133], hl_row=3, row_h=0.62, header_h=0.5)

# 24
s = new()
by = head(s, "에이전트는 이렇게 한 바퀴를 돕니다", "ReAct — Reasoning(추론)과 Acting(행동)을 번갈아 한다는 뜻.")
fb = foot(s, [("고리가 닫혔다. 4일 뒤 여러분이 만들 것이 정확히 이 아홉 단계다.", True)])
table_draw(s, by + Inches(0.05),
    ["단계", "무엇을 하는가"],
    [["01 생각", "온도가 이상한지 보려면 최근 값을 알아야겠다"],
     ["02 행동", "온도 조회 도구를 부른다"],
     ["03 관찰", "계속 오른다"],
     ["04 생각", "네 시간째 한 번도 안 내려갔다. 정상 변동이 아니다"],
     ["05 행동", "정비 이력 조회 도구를 부른다"],
     ["06 관찰", "마지막 정비가 90일 전이다"],
     ["07 생각", "냉각 쪽 문제로 보인다. 속도를 낮추자"],
     ["08 행동", "속도 조절 도구를 부른다"],
     ["09 관찰", "온도가 내려가기 시작한다"]],
    [1.9, 9.833], row_h=0.415, header_h=0.42, size=14)

# 25
s = new()
by = head(s, "하나에게 다 시키면 어떻게 될까요", "앞 장의 아홉 단계를 지시 하나에 몰아넣으면 벌어지는 일.")
fb = foot(s, [("일을 나누고, 나눈 것을 어떻게 이어 붙일지 정해 두는 것 — 이것이 에이전틱 디자인 패턴이다", False)])
pic_fit(s, ART/"art_25.png", int(ML), int(by), int(Inches(4.6)), int(fb - by))
tb(s, Inches(5.8), by + Inches(0.2), Inches(6.7), Inches(1.05),
   '"설비 상태를 점검하고, 이상하면 원인을 찾고, 정비 이력도 보고, 조치안을 쓰고, 급하면 알리고, 보고서 형식은 이러이러하게…"',
   T_BODY, False, GRAY)
body_paras(s, Inches(5.8), Inches(6.7), by + Inches(1.5),
    ["지시가 길어질수록 뒤쪽 요구를 흘린다. 다섯을 시키면 셋만 하고 끝난다.",
     "어디서 틀렸는지 짚을 수가 없다. 통째로 한 번에 하니 중간이 안 보인다.",
     "하나만 고치고 싶어도 전체를 다시 시켜야 한다."])

# 26
s = new()
by = head(s, "하나 — 프롬프트 체이닝",
          ["Prompt Chaining. 프롬프트는 AI에게 주는 지시문, 체이닝은 사슬처럼 잇는다는 뜻.",
           "지시를 여러 도막으로 쪼개 줄줄이 잇는 것이다."])
fb = foot(s, [("각 단계가 짧아 흘리지 않는다. 그리고 어느 단계에서 틀렸는지 바로 보인다. 03이 이상하면 03만 고친다.", False),
              ("조립 라인과 같다. 한 사람이 완제품을 다 만들지 않고 공정마다 한 가지씩 붙여 넘긴다.", False)])
rail_rows(s, by + Inches(0.1), fb - Inches(0.05),
    [("01", "지난 8시간 센서 값에서 이상 구간을 찾아라"),
     ("02", "그 구간의 설비 정비 이력을 조회하라"),
     ("03", "앞의 둘을 근거로 원인을 추정하라"),
     ("04", "조치안을 작성하라")],
    icon_file="art_26.png", desc_size=16)

# 27
s = new()
by = head(s, "둘 — 라우팅", "Routing. 들어온 요청을 살펴보고, 그 일을 가장 잘 처리할 쪽으로 보낸다.")
fb = foot(s, [('같은 "온도 이상"이라도 등급에 따라 가는 길이 다르다. 전부 똑같이 처리하면 사소한 것에 매번 사람을 부르거나, 심각한 것을 기록만 하고 넘긴다.', False),
              ("교차로의 교통경찰과 같다. 차를 직접 몰지 않고 어느 길로 갈지만 정해 준다.", False)])
rail_rows(s, by + Inches(0.15), fb - Inches(0.1),
    [("심각", "즉시 설비를 세우고 담당자에게 전화까지 건다"),
     ("경고", "속도를 낮추고 정비 작업지시를 발행한다"),
     ("주의", "기록만 남기고 교대 리포트에 모아 둔다")],
    icon_file="art_27.png", head_colors=[ORANGE, NAVY, NAVY], desc_size=16)

# 28
s = new()
by = head(s, "셋 — 병렬 처리", "Parallelization. 서로 영향을 주지 않는 일은 동시에 시킨다.")
fb = foot(s, [("여섯 배가 걸리던 일이, 가장 오래 걸린 한 대만큼으로 줄어듭니다.", True)])
pic_fit(s, ART/"art_28.png", int(ML), int(by), int(Inches(5.6)), int(fb - by))
body_paras(s, Inches(6.8), Inches(5.7), by + Inches(0.3),
    ["야간 점검에서 CNC 6대를 한 대씩 차례로 볼 이유가 없다. EQ-01부터 EQ-06까지 동시에 점검하고 끝난 것부터 결과를 모은다.",
     "차례로 보면 여섯 배가 걸리지만, 동시에 보면 가장 오래 걸린 한 대만큼만 걸린다.",
     '다만 순서가 중요한 일에는 못 쓴다. "이상을 찾고 → 그 설비의 이력을 본다"는 앞이 끝나야 뒤가 가능하다.'], gap=0.1)

# 29
s = new()
by = head(s, "넷 — 총괄이 나누고 담당이 처리한다",
          "Orchestrator-Worker. 오케스트레이터는 지휘자, 워커는 실제로 일하는 담당이다. 지휘자는 직접 일하지 않는다.")
fb = foot(s, [("총괄이 셋의 결과를 묶어 하나의 보고서로 낸다.", False),
              ("이것이 넷째 날 여러분이 만들 구조다.", True)])
pic_fit(s, ART/"art_29.png", int(ML), int(by + Inches(0.05)), int(CW), int(fb - by - Inches(0.1)))

# 30
s = new()
by = head(s, "다섯 — 검사하는 쪽을 따로 둔다",
          "Evaluator-Optimizer. 만드는 쪽과 검사하는 쪽을 나눠 두면, 기준에 못 미칠 때 그 단계만 다시 시킬 수 있다.")
fb = foot(s, [("하나라도 빠지면 처음부터 다시 시키지 않는다. 빠진 그 단계만 다시 시킨다.", False),
              ("AI는 그럴듯한 문장을 잘 만든다. 그럴듯한 것과 맞는 것은 다르다. 앞에서 말한 환각을 걸러 내는 자리이기도 하다.", False),
              ("품질 검사 공정과 같다. 만든 사람이 자기 것을 검사하지 않는다.", True)])
pic_fit(s, ART/"art_30.png", int(ML), int(by + Inches(0.05)), int(CW), int(fb - by - Inches(0.1)))

# 31
s = new()
by = head(s, "다섯을 한 줄로", "외울 필요는 없다. 4일 동안 하나씩 직접 쓰게 된다.")
fb = foot(s, [("오늘은 이름만 알면 된다. 넷째 날, 이 중 넷을 여러분 손으로 짠다.", True)])
table_draw(s, by + Inches(0.1),
    ["패턴", "한 줄 뜻", "공장에서", "언제 만나나"],
    [["프롬프트 체이닝", "쪼개서 줄줄이 잇는다", "이상 탐색 → 이력 조회 → 원인 추정 → 조치안", "Day 3"],
     ["라우팅", "성격을 보고 알맞은 곳으로 보낸다", "알람 등급별 대응 분기", "Day 4 확장 미션"],
     ["병렬 처리", "상관없는 일은 동시에", "CNC 6대 동시 점검", "Day 4 확장 미션"],
     ["총괄과 담당", "총괄이 나누고 담당이 처리한다", "감지·진단·조치 담당 + 총괄", "Day 4 기본 미션"],
     ["검증과 되짚기", "검사하는 쪽을 따로 둔다", "근거 없는 리포트를 돌려보낸다", "Day 3~4"]],
    [2.2, 3.3, 4.2, 2.033], hl_row=3, row_h=0.62, header_h=0.5, size=14)

# 32
s = new()
by = head(s, "오늘의 무대 — K-정밀", "경남 창원의 정밀기계 부품 제조사. 다만 이 공장은 전부 소프트웨어입니다.")
fb = foot(s, [('이 기계들은 전부 가짜다. 데이터베이스 테이블의 행 8개이고, "설비가 돈다"는 것은 1초마다 숫자를 바꿔 쓰는 프로그램이다.', True)])
pic_fit(s, ART/"art_32.png", int(ML), int(by), int(Inches(5.2)), int(fb - by))
body_paras(s, Inches(6.4), Inches(6.1), by + Inches(0.25),
    ["AMR은 Autonomous Mobile Robot, 자율이동로봇이다. 정해진 레일 위만 다니는 기존 운반 장비와 달리 스스로 경로를 찾아 공장 안을 돌아다닌다.",
     "읽는 값 — CNC는 온도·진동·rpm·가동상태, AMR은 위치·배터리·적재상태. 이 값들이 1초마다 바뀌며 데이터베이스에 쌓인다.",
     "온도와 진동은 기계가 지금 무리하고 있는가를 알려 주는 대표 신호라 셋째 날 이상감지의 재료가 된다."], gap=0.1)

# 33  — 실화면 캡처 사용
s = new()
by = head(s, "지금, 이 공장은 돌고 있습니다", "강사 화면으로 실제 시뮬레이터를 열어 보여 주는 자리.")
fb = foot(s, [("그 네 개는 오늘 열지 않는다. 넷째 날에 연다.", True),
              ("그리고 여러분 각자에게 자기 몫의 공장이 따로 주어진다. 옆 사람이 내린 명령이 내 공장에 영향을 주지 않는다.", False)])
rail_rows(s, by + Inches(0.1), fb - Inches(0.05),
    [("하나", "여섯 대의 온도·진동·rpm 숫자가 1초마다 바뀐다. 아무도 손대지 않는데 계속 움직인다."),
     ("둘", "강사 콘솔에서 이상을 일부러 만들 수 있다. 온도가 서서히 오르는 드리프트, 진동이 순간 튀는 스파이크, 센서 값이 아예 안 들어오는 결측 — 세 가지다."),
     ("셋", "명령을 넣는 통로가 네 개 준비돼 있다. 설비 속도 조절, 설비 정지, 로봇 파견, 알람 확인 처리.")],
    x0=ML, x1=Inches(6.7), desc_size=14)
pic_fit(s, CAP/"2D공장뷰_Day1시연.png", int(Inches(7.0)), int(by + Inches(0.15)), int(Inches(5.5)), int(fb - by - Inches(0.3)))

# 34
s = new()
by = head(s, "4일 뒤, 이 공장을 움직이는 것은 여러분의 에이전트입니다",
          "오늘 보고, 내일 설계하고, 셋째 날 알아채고, 넷째 날 움직입니다.", title_w=CW)
fb = foot(s, [("그때 첫날 그린 고리가 닫힌다. 인지 → 판단 → 행동 → 다시 인지.", False)])
pic_fit(s, ART/"art_34.png", int(ML), int(by), int(Inches(5.1)), int(fb - by))
body_paras(s, Inches(6.3), Inches(6.2), by + Inches(0.4),
    ["강사가 콘솔에서 EQ-03에 온도 드리프트를 주입한다. 새벽 3시에 벌어졌던 그 일이다.",
     "여러분이 만든 감지 에이전트가 포착하고, 진단 에이전트가 원인을 추정하고, 조치 에이전트가 속도를 낮추면서 정비 로봇 파견 승인을 요청한다.",
     "사람이 승인을 누르면 — 화면 위에서 AMR이 실제로 움직인다."], gap=0.12)

# 35
s = new()
by = head(s, "오늘은 이 공장을 「보는」 눈을 만듭니다", "코드는 한 줄도 치지 않습니다.")
fb = foot(s, [("말로 시키는 개발이 어디까지 되는지, 그리고 어디서부터 안 되는지도 함께 확인한다. 안 되는 지점이 내일 배울 내용이다.", True)])
pic_fit(s, ART/"art_35.png", int(ML), int(by), int(Inches(5.1)), int(fb - by))
tb(s, Inches(6.3), by + Inches(0.05), Inches(6.2), Inches(0.6), "바이브코딩", T_WORD, True, NAVY)
tb(s, Inches(6.3), by + Inches(0.7), Inches(6.2), Inches(0.35), "코드를 직접 쓰지 않고 말로 시켜 만드는 방식", T_WSUB, False, GRAY)
body_paras(s, Inches(6.3), Inches(6.2), by + Inches(1.25),
    ["폐루프의 첫 칸은 인지다. 판단도 행동도 보는 것부터 시작한다. 그래서 오늘 만드는 것은 관제 대시보드다.",
     "대시보드란 흩어진 값을 한 화면에 모아 놓고 계속 갱신해 주는 화면이다. 자동차 계기판과 같은 말이다.",
     "파이썬을 처음 봐도 상관없다. 오늘 하는 일은 무엇을 만들지 정확히 말하는 것이다."], gap=0.08)

# 36
s = new()
by = head(s, "오늘 실습은 네 도막입니다", "Lab 1-1 · 110분.")
fb = foot(s, [('"차이가 뭘까요?"', True)])
rail_rows(s, by + Inches(0.08), fb - Inches(0.05),
    [("Step 1 · 40분 · 작은 성공", "미리 준비된 프롬프트를 시작점으로 받는다. 시뮬레이터에서 값을 읽어 CNC 6대와 AMR 2대의 상태를 보여 주는 실시간 대시보드가 만들어진다. 브라우저를 열면 자기 공장이 움직이는 게 보인다. 그다음 자유롭게 한 가지를 자기 마음대로 바꿔 본다."),
     ("Step 2 · 40분 · 요구사항 추가", "요구사항이 적힌 카드 5장을 순서대로 하나씩 받는다. 온도 임계 알람 표시 → 설비별 이력 그래프 → AMR 경로 표시 → 알람 담당자 배정 → 야간모드·권한 분리."),
     ("Step 3 · 20분 · 강사 시연", "강사가 똑같은 요구사항을 다른 방식으로 지시하는 것을 보여 준다. 여기서는 질문 하나만 던진다."),
     ("Step 4 · 10분 · 남기기", "저장소에 커밋하고 그동안 쓴 프롬프트 기록을 저장한다.")],
    head_size=16, desc_size=13.5)

# 37
s = new()
by = head(s, "시작 전에 30분, 손부터 풉니다", "셋업 랩 — 개발환경 확인과 파이썬 브릿지.")
fb = foot(s, [("여러분은 이미 더 어려운 걸 배웠습니다.", True)])
rail_rows(s, by + Inches(0.1), fb - Inches(0.05),
    [("01", "VS Code에 Claude 플러그인이 제대로 깔렸는지 확인하고, GitHub에 로그인해 실습 저장소를 내려받는다. 설치 안내는 미리 문서로 보냈으므로 현장에서는 안 되는 사람을 고쳐 주는 데 시간을 쓴다."),
     ("02", "시뮬레이터에 접속되는지 확인한다. 자기 몫으로 배정된 공장 화면을 열어 본다."),
     ("03", "파이썬 브릿지 한 장. C를 배운 사람이 파이썬으로 넘어올 때 헷갈리는 것들을 한 장짜리 대응표로 준다. 포인터는 참조로, for문·함수·클래스는 이렇게 대응된다는 식이다.")],
    desc_size=14.5)

# 38
s = new()
cover(s, "Day 1 · 보기 → Day 2 · 인지의 구조화",
      ["오늘 공장을 봤습니다.", "내일은 공장을 설계합니다."], "",
      [("오늘 남는 것 — 관제 화면 하나. 그리고 말로 시켜 만들 때 어디까지 되고 어디서부터 안 되는지에 대한 감각.", False),
       ("랩업 20분에는 지원자 두세 명이 자기 대시보드를 공유하고, 누가 가장 창의적으로 바꿨는지 그 자리에서 투표한다.", False),
       ("내일 첫 10분에 어제 왜 그렇게 됐는지를 공개한다.", False)])

prs.save(OUT)
print(f"저장: {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)}장)")
